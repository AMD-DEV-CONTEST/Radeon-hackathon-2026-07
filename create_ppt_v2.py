from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色主题 - AMD风格
AMD_RED = RGBColor(0xED, 0x1C, 0x24)
DARK_BG = RGBColor(0x1F, 0x1F, 0x1F)
LIGHT_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_GRAY = RGBColor(0xC8, 0xD3, 0xE6)
TEAL = RGBColor(0x0D, 0x94, 0x88)
GREEN = RGBColor(0x2C, 0x5F, 0x2D)
BLUE = RGBColor(0x12, 0x45, 0x7A)
ORANGE = RGBColor(0xF0, 0xA2, 0x02)

# ========== 幻灯片1：标题页 ==========
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

# 背景
bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BG
bg.line.fill.background()

# AMD红色装饰条
accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), prs.slide_width, Inches(0.1))
accent.fill.solid()
accent.fill.fore_color.rgb = AMD_RED
accent.line.fill.background()

# 标题
title = slide1.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11.3), Inches(1.2))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "Supply Chain DocAgent"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(48)
r.font.bold = True
r.font.color.rgb = LIGHT_TEXT

# 副标题
subtitle = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(0.8))
tf = subtitle.text_frame
p = tf.paragraphs[0]
p.text = "基于 AMD Radeon GPU 的供应链单据智能处理 Agent"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(24)
r.font.color.rgb = ACCENT_GRAY

# 团队信息
team = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.3), Inches(1))
tf = team.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "AMD AI DevMaster Hackathon 2026 | Track 2: Private AI Agent Development"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(18)
r.font.color.rgb = TEAL

# 技术栈
tech = slide1.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.3), Inches(1))
tf = tech.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Python | LangChain | vLLM | ROCm | Llama-3.1-8B | BGE Embeddings | Gradio"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(14)
r.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

# ========== 幻灯片2：问题与痛点 ==========
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "供应链单据处理的痛点"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# 痛点表格
pain_points = [
    ("痛点", "现状", "影响"),
    ("人工逐单核对", "每张单据处理需 15 分钟", "效率低下，难以应对大批量"),
    ("三单匹配依赖经验", "新手容易漏检差异", "错误率约 3%"),
    ("单据格式不统一", "不同供应商模板各异", "识别困难，需人工适配"),
    ("异常处理追溯难", "纸质单据难以检索", "审计困难，责任不清"),
]

for row_idx, row in enumerate(pain_points):
    y_pos = Inches(1.3 + row_idx * 1.1)
    is_header = row_idx == 0
    
    for col_idx, cell_text in enumerate(row):
        x_pos = Inches(0.5 + col_idx * 4.2)
        width = Inches(4)
        
        box = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y_pos, width, Inches(0.9))
        
        if is_header:
            box.fill.solid()
            box.fill.fore_color.rgb = DARK_BG
        else:
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5) if row_idx % 2 == 0 else LIGHT_TEXT
        box.line.fill.background()
        
        tb = slide2.shapes.add_textbox(x_pos + Inches(0.2), y_pos + Inches(0.2), width - Inches(0.4), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cell_text
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(14) if not is_header else Pt(16)
        p.runs[0].font.bold = is_header
        p.runs[0].font.color.rgb = LIGHT_TEXT if is_header else DARK_BG

# ========== 幻灯片3：解决方案概述 ==========
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "我们的解决方案"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# 核心特性
features = [
    ("零云端依赖", "所有数据本地处理\n保护商业机密", GREEN),
    ("全流程自动化", "识别→提取→校验\n→录入→异常处理", BLUE),
    ("多Agent协作", "5个专业Agent\n各司其职", TEAL),
    ("人机协作", "AI处理标准流程\n异常自动升级人工", ORANGE),
]

for i, (title_text, desc, color) in enumerate(features):
    x_pos = Inches(0.5 + i * 3.2)
    
    box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.5), Inches(2.8), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    
    tb = slide3.shapes.add_textbox(x_pos + Inches(0.2), Inches(1.8), Inches(2.4), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(22)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = LIGHT_TEXT
    
    tb2 = slide3.shapes.add_textbox(x_pos + Inches(0.2), Inches(2.6), Inches(2.4), Inches(1.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.alignment = PP_ALIGN.CENTER
    p2.runs[0].font.size = Pt(16)
    p2.runs[0].font.color.rgb = LIGHT_TEXT

# 技术亮点
note = slide3.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(2.5))
tf = note.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "技术亮点"
p.runs[0].font.size = Pt(24)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

highlights = [
    "• Llama-3.1-8B + ROCm 6.x 本地推理，FP16 混合精度",
    "• BGE Embeddings 向量化 + LlamaIndex RAG 引擎",
    "• vLLM 高性能推理服务，GPU 显存占用 < 12GB",
    "• Gradio Web UI，支持对话式查询和文档处理",
]

for i, highlight in enumerate(highlights):
    p2 = tf.add_paragraph()
    p2.text = highlight
    p2.runs[0].font.size = Pt(16)
    p2.runs[0].font.color.rgb = RGBColor(0x44, 0x44, 0x44)

# ========== 幻灯片4：系统架构 ==========
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "系统架构"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# 架构图
layers = [
    ("输入层", "拍照/扫描 | PDF/Excel | 邮件附件", RGBColor(0xE9, 0xDF, 0xC7)),
    ("Agent 编排层", "Classify | Extract | Validate | Entry | Exception", TEAL),
    ("RAG 引擎", "BGE Embeddings + LlamaIndex | 单据模板库 + 历史案例库", BLUE),
    ("LLM 推理", "Llama-3.1-8B + vLLM + ROCm | FP16 混合精度", AMD_RED),
    ("输出层", "核对通过 → 自动入库 | 异常 → 推送审批", GREEN),
]

for i, (title_text, desc, color) in enumerate(layers):
    y_pos = Inches(1.2 + i * 1.1)
    
    box = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), y_pos, Inches(11), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    
    tb = slide4.shapes.add_textbox(Inches(1.5), y_pos + Inches(0.1), Inches(3), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.runs[0].font.size = Pt(20)
    p.runs[0].font.bold = True
    
    tb2 = slide4.shapes.add_textbox(Inches(5), y_pos + Inches(0.1), Inches(6.5), Inches(0.7))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.runs[0].font.size = Pt(14)

# ========== 幻灯片5：Agent 分工详情 ==========
slide5 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "多 Agent 协作分工"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# Agent表格
agents = [
    ("Agent", "职责", "输入", "输出"),
    ("Classify Agent", "识别单据类型", "单据文本", "类型标签 (PO/送货单/发票)"),
    ("Extract Agent", "提取结构化字段", "单据文本 + 类型", "JSON 字段"),
    ("Validate Agent", "三单交叉校验", "PO + 送货单 + 发票", "校验结果 + 差异"),
    ("Entry Agent", "ERP 数据回写", "校验通过的数据", "录入确认"),
    ("Exception Agent", "异常分类与通知", "校验失败的数据", "审批通知"),
]

for row_idx, row in enumerate(agents):
    y_pos = Inches(1.2 + row_idx * 0.9)
    is_header = row_idx == 0
    
    col_widths = [Inches(2.5), Inches(3), Inches(3.5), Inches(3.5)]
    x_pos = Inches(0.5)
    
    for col_idx, (cell_text, col_width) in enumerate(zip(row, col_widths)):
        box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y_pos, col_width, Inches(0.8))
        
        if is_header:
            box.fill.solid()
            box.fill.fore_color.rgb = DARK_BG
        else:
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5) if row_idx % 2 == 0 else LIGHT_TEXT
        box.line.fill.background()
        
        tb = slide5.shapes.add_textbox(x_pos + Inches(0.15), y_pos + Inches(0.15), col_width - Inches(0.3), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cell_text
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(13) if not is_header else Pt(14)
        p.runs[0].font.bold = is_header
        p.runs[0].font.color.rgb = LIGHT_TEXT if is_header else DARK_BG
        
        x_pos += col_width

# ========== 幻灯片6：技术栈详情 ==========
slide6 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "技术栈详情"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# 技术栈表格
tech_stack = [
    ("组件", "技术选型", "版本/配置", "作用"),
    ("LLM", "Llama-3.1-8B-Instruct", "FP16, temp=0.3", "单据理解与生成"),
    ("推理引擎", "vLLM + ROCm 6.x", "批处理优化", "高性能 GPU 推理"),
    ("Embeddings", "BGE-small-en-v1.5", "512 chunk, top_k=5", "向量化与 RAG 检索"),
    ("RAG 框架", "LlamaIndex", "向量数据库", "知识库管理"),
    ("Agent 框架", "LangChain", "多Agent编排", "工作流协调"),
    ("Web UI", "Gradio", "端口 7860", "用户交互界面"),
]

for row_idx, row in enumerate(tech_stack):
    y_pos = Inches(1.2 + row_idx * 0.8)
    is_header = row_idx == 0
    
    col_widths = [Inches(2), Inches(3.5), Inches(3), Inches(3.5)]
    x_pos = Inches(0.5)
    
    for col_idx, (cell_text, col_width) in enumerate(zip(row, col_widths)):
        box = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y_pos, col_width, Inches(0.7))
        
        if is_header:
            box.fill.solid()
            box.fill.fore_color.rgb = DARK_BG
        else:
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5) if row_idx % 2 == 0 else LIGHT_TEXT
        box.line.fill.background()
        
        tb = slide6.shapes.add_textbox(x_pos + Inches(0.1), y_pos + Inches(0.1), col_width - Inches(0.2), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cell_text
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(12) if not is_header else Pt(13)
        p.runs[0].font.bold = is_header
        p.runs[0].font.color.rgb = LIGHT_TEXT if is_header else DARK_BG
        
        x_pos += col_width

# ========== 幻灯片7：AMD GPU 优化 ==========
slide7 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "AMD Radeon GPU + ROCm 优化"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# 测试环境
env_title = slide7.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6), Inches(0.5))
tf = env_title.text_frame
p = tf.paragraphs[0]
p.text = "测试环境"
p.runs[0].font.size = Pt(24)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

env_items = [
    "• GPU: AMD Radeon RX 7900 XTX (24GB VRAM)",
    "• OS: Ubuntu 22.04 LTS",
    "• ROCm: 6.1",
    "• PyTorch: 2.4.0+rocm6.1",
]

env_box = slide7.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6), Inches(2.5))
tf = env_box.text_frame
tf.word_wrap = True
for i, item in enumerate(env_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = item
    p.runs[0].font.size = Pt(16)

# 性能指标
perf_title = slide7.shapes.add_textbox(Inches(7), Inches(1.2), Inches(6), Inches(0.5))
tf = perf_title.text_frame
p = tf.paragraphs[0]
p.text = "性能指标"
p.runs[0].font.size = Pt(24)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

perf_items = [
    ("单据分类准确率", "98%+"),
    ("字段提取准确率", "95%+"),
    ("三单匹配准确率", "99%+"),
    ("单张处理时间", "< 30 秒"),
    ("GPU 显存占用", "< 12 GB"),
]

perf_box = slide7.shapes.add_textbox(Inches(7), Inches(1.8), Inches(6), Inches(3))
tf = perf_box.text_frame
tf.word_wrap = True
for i, (metric, value) in enumerate(perf_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"{metric}: {value}"
    p.runs[0].font.size = Pt(16)

# 优化技术
opt_title = slide7.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12), Inches(0.5))
tf = opt_title.text_frame
p = tf.paragraphs[0]
p.text = "优化技术"
p.runs[0].font.size = Pt(24)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

opt_items = [
    "• FP16 混合精度推理 - 减少显存占用，提升计算效率",
    "• vLLM 批处理优化 - 支持并发请求，提高吞吐量",
    "• BGE 向量化 - 本地 Embedding 计算，零网络延迟",
    "• ROCm HIP 后端 - 原生 AMD GPU 支持，无需 CUDA",
]

opt_box = slide7.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(12), Inches(2))
tf = opt_box.text_frame
tf.word_wrap = True
for i, item in enumerate(opt_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = item
    p.runs[0].font.size = Pt(14)
    p.runs[0].font.color.rgb = RGBColor(0x44, 0x44, 0x44)

# ========== 幻灯片8：三单校验逻辑 ==========
slide8 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "三单交叉校验逻辑"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# 校验规则
rules_title = slide8.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6), Inches(0.5))
tf = rules_title.text_frame
p = tf.paragraphs[0]
p.text = "校验规则"
p.runs[0].font.size = Pt(24)
p.runs[0].font.bold = True

rules = [
    ("PO号匹配", "po_number ↔ po_reference"),
    ("物料编码匹配", "item_code ↔ item_code"),
    ("数量匹配", "quantity ↔ quantity_delivered (±5%容差)"),
    ("单价匹配", "unit_price ↔ unit_price"),
    ("金额匹配", "total_amount ↔ amount (±1%容差)"),
]

rules_box = slide8.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6), Inches(3.5))
tf = rules_box.text_frame
tf.word_wrap = True
for i, (rule, desc) in enumerate(rules):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"• {rule}"
    p.runs[0].font.size = Pt(16)
    p.runs[0].font.bold = True
    p2 = tf.add_paragraph()
    p2.text = f"  {desc}"
    p2.runs[0].font.size = Pt(14)
    p2.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# 校验流程
flow_title = slide8.shapes.add_textbox(Inches(7), Inches(1.2), Inches(6), Inches(0.5))
tf = flow_title.text_frame
p = tf.paragraphs[0]
p.text = "校验流程"
p.runs[0].font.size = Pt(24)
p.runs[0].font.bold = True

flow_steps = [
    ("1. 上传采购订单", "Extract Agent 提取字段"),
    ("2. 上传送货单", "Extract Agent 提取字段"),
    ("3. 上传发票", "Extract Agent 提取字段"),
    ("4. 三单交叉校验", "Validate Agent 比对"),
    ("5. 生成校验报告", "通过/异常 分流处理"),
]

flow_box = slide8.shapes.add_textbox(Inches(7), Inches(1.8), Inches(6), Inches(3.5))
tf = flow_box.text_frame
tf.word_wrap = True
for i, (step, desc) in enumerate(flow_steps):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"{step}"
    p.runs[0].font.size = Pt(16)
    p.runs[0].font.bold = True
    p2 = tf.add_paragraph()
    p2.text = f"  {desc}"
    p2.runs[0].font.size = Pt(14)
    p2.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# ========== 幻灯片9：商业价值 ==========
slide9 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide9.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "商业价值"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# 价值表格
value_data = [
    ("指标", "改进前", "改进后", "提升"),
    ("单据处理时间", "15 分钟/张", "2 分钟/张", "87%"),
    ("错误率", "3%", "0.5%", "83%"),
    ("人力需求", "8 人", "2 人（异常处理）", "75%"),
    ("库存差异发现时效", "T+1", "实时", "100%"),
]

for row_idx, row in enumerate(value_data):
    y_pos = Inches(1.3 + row_idx * 1.1)
    is_header = row_idx == 0
    
    col_widths = [Inches(3.5), Inches(3), Inches(3), Inches(2.5)]
    x_pos = Inches(0.5)
    
    for col_idx, (cell_text, col_width) in enumerate(zip(row, col_widths)):
        box = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y_pos, col_width, Inches(0.9))
        
        if is_header:
            box.fill.solid()
            box.fill.fore_color.rgb = DARK_BG
        elif col_idx == 3:  # 提升列用绿色
            box.fill.solid()
            box.fill.fore_color.rgb = GREEN
        else:
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5) if row_idx % 2 == 0 else LIGHT_TEXT
        box.line.fill.background()
        
        tb = slide9.shapes.add_textbox(x_pos + Inches(0.1), y_pos + Inches(0.2), col_width - Inches(0.2), Inches(0.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = cell_text
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(16) if not is_header else Pt(18)
        p.runs[0].font.bold = is_header or col_idx == 3
        p.runs[0].font.color.rgb = LIGHT_TEXT if is_header or col_idx == 3 else DARK_BG
        
        x_pos += col_width

# ========== 幻灯片10：创新点 ==========
slide10 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide10.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "创新点"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# 创新点
innovations = [
    ("领域专精 Agent", "针对供应链单据场景深度优化\n非通用 Agent，专业能力更强", TEAL),
    ("多 Agent 协作", "5 个专业 Agent 各司其职\n协同完成复杂业务流程", BLUE),
    ("三单交叉校验", "自动核对 PO/送货单/发票\n替代人工经验判断", GREEN),
    ("完全本地化", "商业数据零泄露风险\n适合企业私有化部署", AMD_RED),
]

for i, (title_text, desc, color) in enumerate(innovations):
    x_pos = Inches(0.5 + i * 3.2)
    
    box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.5), Inches(2.8), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    
    tb = slide10.shapes.add_textbox(x_pos + Inches(0.2), Inches(1.8), Inches(2.4), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(22)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = LIGHT_TEXT
    
    tb2 = slide10.shapes.add_textbox(x_pos + Inches(0.2), Inches(2.6), Inches(2.4), Inches(1.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.alignment = PP_ALIGN.CENTER
    p2.runs[0].font.size = Pt(16)
    p2.runs[0].font.color.rgb = LIGHT_TEXT

# ========== 幻灯片11：未来扩展 ==========
slide11 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide11.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "未来扩展"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# 扩展计划
expansions = [
    ("图像单据支持", "集成多模态 LLM\n支持拍照识别手写单据"),
    ("全流程覆盖", "扩展到出库、对账、结算\n覆盖供应链全链路"),
    ("ERP 集成", "对接 SAP/Oracle/用友\n实现真正的自动化"),
    ("多仓库协同", "支持多仓库并行处理\n统一管理和调度"),
]

for i, (title_text, desc) in enumerate(expansions):
    x_pos = Inches(0.5 + i * 3.2)
    
    box = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, Inches(1.5), Inches(2.8), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    box.line.fill.background()
    
    # 左侧装饰条
    bar = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, Inches(1.5), Inches(0.1), Inches(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = [TEAL, BLUE, GREEN, ORANGE][i]
    bar.line.fill.background()
    
    tb = slide11.shapes.add_textbox(x_pos + Inches(0.3), Inches(1.8), Inches(2.3), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.runs[0].font.size = Pt(20)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = DARK_BG
    
    tb2 = slide11.shapes.add_textbox(x_pos + Inches(0.3), Inches(2.6), Inches(2.3), Inches(1.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.runs[0].font.size = Pt(14)
    p2.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# ========== 幻灯片12：总结 ==========
slide12 = prs.slides.add_slide(prs.slide_layouts[6])

# 背景
bg = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BG
bg.line.fill.background()

# 标题
title = slide12.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11.3), Inches(1.2))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(56)
r.font.bold = True
r.font.color.rgb = LIGHT_TEXT

# 关键特性
features = [
    "✓ Multi-Agent collaboration (5 specialized agents)",
    "✓ AMD GPU acceleration with ROCm 6.x",
    "✓ Complete privacy protection (100% local)",
    "✓ Real-world business value (87% time reduction)",
    "✓ High accuracy (98%+ classification, 99%+ matching)",
]

for i, feature in enumerate(features):
    tb = slide12.shapes.add_textbox(Inches(2), Inches(2.8 + i * 0.5), Inches(9), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = feature
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(20)
    p.runs[0].font.color.rgb = TEAL

# 联系信息
contact = slide12.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11.3), Inches(1.2))
tf = contact.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Team: cyslmsolomon"
p.alignment = PP_ALIGN.CENTER
p.runs[0].font.size = Pt(18)
p.runs[0].font.color.rgb = ACCENT_GRAY

p2 = tf.add_paragraph()
p2.text = "GitHub: github.com/cyslmsolomon/Radeon-hackathon-2026-07"
p2.alignment = PP_ALIGN.CENTER
p2.runs[0].font.size = Pt(16)
p2.runs[0].font.color.rgb = RGBColor(0x88, 0x88, 0x88)

# 保存演示文稿
output_path = "Supply_Chain_DocAgent_Presentation_v2.pptx"
prs.save(output_path)
print(f"演示文稿已保存: {output_path}")
