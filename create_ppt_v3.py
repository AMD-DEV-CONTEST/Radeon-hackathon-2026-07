from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ==================== 专业配色方案 ====================
# AMD品牌色 + 科技感专业配色
# 主色: 深邃科技蓝 (#0D1B2A) - 用于背景和标题
# 辅助色: 银灰 (#E0E0E0) - 用于正文
# 强调色: AMD红 (#ED1C24) - 用于重点强调
# 点缀色: 科技青 (#00B4D8) - 用于数据和链接
# 背景色: 渐变深蓝 (#0D1B2A → #1B263B)

# 颜色定义
DEEP_BLUE = RGBColor(0x0D, 0x1B, 0x2A)      # 深邃科技蓝（主背景）
DARK_NAVY = RGBColor(0x1B, 0x26, 0x3B)      # 深海军蓝（次背景）
MID_BLUE = RGBColor(0x25, 0x3A, 0x54)       # 中蓝（卡片背景）
LIGHT_BLUE = RGBColor(0x41, 0x5A, 0x77)     # 浅蓝灰

WHITE = RGBColor(0xFF, 0xFF, 0xFF)           # 纯白
OFF_WHITE = RGBColor(0xF0, 0xF0, 0xF0)      # 柔和白
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)     # 浅灰（正文）
MID_GRAY = RGBColor(0xA0, 0xA0, 0xA0)       # 中灰（辅助文字）
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)      # 深灰

AMD_RED = RGBColor(0xED, 0x1C, 0x24)        # AMD品牌红（强调）
TECH_CYAN = RGBColor(0x00, 0xB4, 0xD8)      # 科技青（数据/链接）
SUCCESS_GREEN = RGBColor(0x2E, 0xCC, 0x71)  # 成功绿
WARNING_ORANGE = RGBColor(0xF3, 0x9C, 0x12) # 警告橙

# 字体定义
FONT_TITLE = "Segoe UI Semibold"
FONT_BODY = "Segoe UI"
FONT_MONO = "Consolas"

def add_background(slide, color=DEEP_BLUE):
    """添加深色背景"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_accent_line(slide, y_pos=Inches(1.1), color=AMD_RED):
    """添加强调线条"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y_pos, Inches(11.7), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def add_title(slide, text, y_pos=Inches(0.3), color=WHITE, size=Pt(36)):
    """添加标题"""
    tb = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(11.7), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.size = size
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = color
    p.runs[0].font.name = FONT_TITLE
    return tb

def add_subtitle(slide, text, y_pos=Inches(1.2), color=LIGHT_GRAY, size=Pt(18)):
    """添加副标题"""
    tb = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(11.7), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.size = size
    p.runs[0].font.color.rgb = color
    p.runs[0].font.name = FONT_BODY
    return tb

def add_card(slide, x, y, w, h, title_text, content_lines, accent_color=TECH_CYAN):
    """添加卡片式内容块"""
    # 卡片背景
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = MID_BLUE
    card.line.fill.background()
    
    # 左侧强调条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    
    # 标题
    tb_title = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.15), w - Inches(0.4), Inches(0.4))
    tf = tb_title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.runs[0].font.size = Pt(16)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE
    p.runs[0].font.name = FONT_TITLE
    
    # 内容
    tb_content = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.55), w - Inches(0.4), h - Inches(0.65))
    tf = tb_content.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.runs[0].font.size = Pt(13)
        p.runs[0].font.color.rgb = LIGHT_GRAY
        p.runs[0].font.name = FONT_BODY
        p.space_after = Pt(6)
    
    return card

def add_stat_callout(slide, x, y, number, label, color=TECH_CYAN):
    """添加统计数据突出显示"""
    # 数字
    tb_num = slide.shapes.add_textbox(x, y, Inches(2.5), Inches(1))
    tf = tb_num.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(48)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = color
    p.runs[0].font.name = FONT_TITLE
    
    # 标签
    tb_label = slide.shapes.add_textbox(x, y + Inches(0.8), Inches(2.5), Inches(0.4))
    tf = tb_label.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(14)
    p.runs[0].font.color.rgb = MID_GRAY
    p.runs[0].font.name = FONT_BODY
    
    return tb_num

# ==================== 幻灯片1：标题页 ====================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide1)

# 顶部装饰条
top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = AMD_RED
top_bar.line.fill.background()

# 主标题
tb = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11.3), Inches(1.2))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Supply Chain DocAgent"
p.alignment = PP_ALIGN.CENTER
p.runs[0].font.size = Pt(54)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = WHITE
p.runs[0].font.name = FONT_TITLE

# 副标题
tb2 = slide1.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.3), Inches(0.8))
tf = tb2.text_frame
p = tf.paragraphs[0]
p.text = "基于 AMD Radeon GPU 的供应链单据智能处理 Agent"
p.alignment = PP_ALIGN.CENTER
p.runs[0].font.size = Pt(24)
p.runs[0].font.color.rgb = LIGHT_GRAY
p.runs[0].font.name = FONT_BODY

# 赛道信息
tb3 = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.3), Inches(0.6))
tf = tb3.text_frame
p = tf.paragraphs[0]
p.text = "AMD AI DevMaster Hackathon 2026  |  Track 2: Private AI Agent Development"
p.alignment = PP_ALIGN.CENTER
p.runs[0].font.size = Pt(18)
p.runs[0].font.color.rgb = TECH_CYAN
p.runs[0].font.name = FONT_BODY

# 技术栈
tb4 = slide1.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11.3), Inches(0.5))
tf = tb4.text_frame
p = tf.paragraphs[0]
p.text = "Python  •  LangChain  •  vLLM  •  ROCm  •  Llama-3.1  •  BGE  •  Gradio"
p.alignment = PP_ALIGN.CENTER
p.runs[0].font.size = Pt(14)
p.runs[0].font.color.rgb = MID_GRAY
p.runs[0].font.name = FONT_MONO

# 底部装饰
bottom_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.42), prs.slide_width, Inches(0.08))
bottom_bar.fill.solid()
bottom_bar.fill.fore_color.rgb = AMD_RED
bottom_bar.line.fill.background()

# ==================== 幻灯片2：问题与痛点 ====================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2)
add_title(slide2, "供应链单据处理的痛点")
add_accent_line(slide2)

# 痛点卡片
pain_points = [
    ("人工逐单核对", ["每张单据处理需 15 分钟", "效率低下，难以应对大批量"], AMD_RED),
    ("三单匹配依赖经验", ["新手容易漏检差异", "错误率约 3%"], WARNING_ORANGE),
    ("单据格式不统一", ["不同供应商模板各异", "识别困难，需人工适配"], TECH_CYAN),
    ("异常处理追溯难", ["纸质单据难以检索", "审计困难，责任不清"], MID_GRAY),
]

for i, (title, lines, color) in enumerate(pain_points):
    x = Inches(0.8 + i * 3.1)
    add_card(slide2, x, Inches(1.8), Inches(2.8), Inches(2.2), title, lines, color)

# 底部数据
add_stat_callout(slide2, Inches(1), Inches(4.8), "15 min", "每张单据处理时间", AMD_RED)
add_stat_callout(slide2, Inches(4), Inches(4.8), "3%", "人工错误率", WARNING_ORANGE)
add_stat_callout(slide2, Inches(7), Inches(4.8), "T+1", "库存差异发现时效", TECH_CYAN)
add_stat_callout(slide2, Inches(10), Inches(4.8), "8人", "处理团队规模", MID_GRAY)

# ==================== 幻灯片3：解决方案 ====================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3)
add_title(slide3, "我们的解决方案")
add_accent_line(slide3)

# 核心特性
features = [
    ("零云端依赖", ["所有数据本地处理", "保护商业机密", "符合数据安全法规"], SUCCESS_GREEN),
    ("全流程自动化", ["识别→提取→校验", "→录入→异常处理", "端到端智能处理"], TECH_CYAN),
    ("多Agent协作", ["5个专业Agent", "各司其职", "协同完成复杂流程"], AMD_RED),
    ("人机协作", ["AI处理标准流程", "异常自动升级人工", "效率与灵活性兼顾"], WARNING_ORANGE),
]

for i, (title, lines, color) in enumerate(features):
    x = Inches(0.8 + i * 3.1)
    add_card(slide3, x, Inches(1.8), Inches(2.8), Inches(2.5), title, lines, color)

# 技术亮点
add_subtitle(slide3, "核心技术亮点", y_pos=Inches(4.8))

highlights = [
    "• Llama-3.1-8B + ROCm 6.x 本地推理，FP16 混合精度",
    "• BGE Embeddings 向量化 + LlamaIndex RAG 引擎",
    "• vLLM 高性能推理服务，GPU 显存占用 < 12GB",
    "• Gradio Web UI，支持对话式查询和文档处理",
]

tb_hl = slide3.shapes.add_textbox(Inches(0.8), Inches(5.3), Inches(11.7), Inches(2))
tf = tb_hl.text_frame
tf.word_wrap = True
for i, hl in enumerate(highlights):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = hl
    p.runs[0].font.size = Pt(15)
    p.runs[0].font.color.rgb = LIGHT_GRAY
    p.runs[0].font.name = FONT_BODY

# ==================== 幻灯片4：系统架构 ====================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4)
add_title(slide4, "系统架构")
add_accent_line(slide4)

# 架构图
layers = [
    ("输入层", "拍照/扫描 | PDF/Excel | 邮件附件", MID_GRAY),
    ("Agent 编排层", "Classify | Extract | Validate | Entry | Exception", TECH_CYAN),
    ("RAG 引擎", "BGE Embeddings + LlamaIndex | 单据模板库 + 历史案例库", SUCCESS_GREEN),
    ("LLM 推理", "Llama-3.1-8B + vLLM + ROCm | FP16 混合精度", AMD_RED),
    ("输出层", "核对通过 → 自动入库 | 异常 → 推送审批", WARNING_ORANGE),
]

for i, (title, desc, color) in enumerate(layers):
    y = Inches(1.5 + i * 1.1)
    
    # 层背景
    layer_bg = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(0.9))
    layer_bg.fill.solid()
    layer_bg.fill.fore_color.rgb = MID_BLUE
    layer_bg.line.fill.background()
    
    # 左侧颜色条
    color_bar = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(0.1), Inches(0.9))
    color_bar.fill.solid()
    color_bar.fill.fore_color.rgb = color
    color_bar.line.fill.background()
    
    # 标题
    tb = slide4.shapes.add_textbox(Inches(1.2), y + Inches(0.15), Inches(3), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.runs[0].font.size = Pt(18)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE
    p.runs[0].font.name = FONT_TITLE
    
    # 描述
    tb2 = slide4.shapes.add_textbox(Inches(4.5), y + Inches(0.15), Inches(7.5), Inches(0.6))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.runs[0].font.size = Pt(14)
    p2.runs[0].font.color.rgb = LIGHT_GRAY
    p2.runs[0].font.name = FONT_BODY

# ==================== 幻灯片5：Agent 分工 ====================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide5)
add_title(slide5, "多 Agent 协作分工")
add_accent_line(slide5)

# Agent表格
agents = [
    ("Classify Agent", "识别单据类型", "单据文本", "类型标签", TECH_CYAN),
    ("Extract Agent", "提取结构化字段", "单据文本 + 类型", "JSON 字段", SUCCESS_GREEN),
    ("Validate Agent", "三单交叉校验", "PO + 送货单 + 发票", "校验结果", WARNING_ORANGE),
    ("Entry Agent", "ERP 数据回写", "校验通过的数据", "录入确认", AMD_RED),
    ("Exception Agent", "异常分类与通知", "校验失败的数据", "审批通知", MID_GRAY),
]

# 表头
headers = ["Agent", "职责", "输入", "输出"]
header_widths = [Inches(2.2), Inches(2.8), Inches(3.5), Inches(3.2)]
x_start = Inches(0.8)

# 表头背景
header_bg = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_start, Inches(1.5), Inches(11.7), Inches(0.6))
header_bg.fill.solid()
header_bg.fill.fore_color.rgb = MID_BLUE
header_bg.line.fill.background()

x_pos = x_start
for i, (header, width) in enumerate(zip(headers, header_widths)):
    tb = slide5.shapes.add_textbox(x_pos + Inches(0.1), Inches(1.55), width - Inches(0.2), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = header
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(14)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = TECH_CYAN
    p.runs[0].font.name = FONT_TITLE
    x_pos += width

# 数据行
for row_idx, (agent, duty, input, output, color) in enumerate(agents):
    y = Inches(2.2 + row_idx * 0.9)
    
    # 行背景
    row_bg = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_start, y, Inches(11.7), Inches(0.8))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = DARK_NAVY if row_idx % 2 == 0 else MID_BLUE
    row_bg.line.fill.background()
    
    # 左侧颜色条
    color_bar = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_start, y, Inches(0.06), Inches(0.8))
    color_bar.fill.solid()
    color_bar.fill.fore_color.rgb = color
    color_bar.line.fill.background()
    
    # 数据
    x_pos = x_start + Inches(0.2)
    for i, (cell, width) in enumerate(zip([agent, duty, input, output], header_widths)):
        tb = slide5.shapes.add_textbox(x_pos, y + Inches(0.15), width - Inches(0.3), Inches(0.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = cell
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(13)
        p.runs[0].font.color.rgb = WHITE if i == 0 else LIGHT_GRAY
        p.runs[0].font.bold = i == 0
        p.runs[0].font.name = FONT_BODY
        x_pos += width

# ==================== 幻灯片6：技术栈 ====================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide6)
add_title(slide6, "技术栈详情")
add_accent_line(slide6)

# 技术栈表格
tech_stack = [
    ("LLM", "Llama-3.1-8B-Instruct", "FP16, temp=0.3", "单据理解与生成", AMD_RED),
    ("推理引擎", "vLLM + ROCm 6.x", "批处理优化", "高性能 GPU 推理", TECH_CYAN),
    ("Embeddings", "BGE-small-en-v1.5", "512 chunk, top_k=5", "向量化与 RAG 检索", SUCCESS_GREEN),
    ("RAG 框架", "LlamaIndex", "向量数据库", "知识库管理", WARNING_ORANGE),
    ("Agent 框架", "LangChain", "多Agent编排", "工作流协调", MID_GRAY),
    ("Web UI", "Gradio", "端口 7860", "用户交互界面", LIGHT_BLUE),
]

# 表头
tech_headers = ["组件", "技术选型", "版本/配置", "作用"]
tech_widths = [Inches(2), Inches(3.2), Inches(3), Inches(3.5)]
x_start = Inches(0.8)

# 表头背景
header_bg = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_start, Inches(1.5), Inches(11.7), Inches(0.6))
header_bg.fill.solid()
header_bg.fill.fore_color.rgb = MID_BLUE
header_bg.line.fill.background()

x_pos = x_start
for header, width in zip(tech_headers, tech_widths):
    tb = slide6.shapes.add_textbox(x_pos + Inches(0.1), Inches(1.55), width - Inches(0.2), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = header
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(14)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = TECH_CYAN
    p.runs[0].font.name = FONT_TITLE
    x_pos += width

# 数据行
for row_idx, (comp, tech, config, role, color) in enumerate(tech_stack):
    y = Inches(2.2 + row_idx * 0.8)
    
    row_bg = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_start, y, Inches(11.7), Inches(0.7))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = DARK_NAVY if row_idx % 2 == 0 else MID_BLUE
    row_bg.line.fill.background()
    
    color_bar = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_start, y, Inches(0.06), Inches(0.7))
    color_bar.fill.solid()
    color_bar.fill.fore_color.rgb = color
    color_bar.line.fill.background()
    
    x_pos = x_start + Inches(0.2)
    for cell, width in zip([comp, tech, config, role], tech_widths):
        tb = slide6.shapes.add_textbox(x_pos, y + Inches(0.12), width - Inches(0.3), Inches(0.45))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = cell
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.name = FONT_BODY
        x_pos += width

# ==================== 幻灯片7：AMD GPU 优化 ====================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide7)
add_title(slide7, "AMD Radeon GPU + ROCm 优化")
add_accent_line(slide7)

# 测试环境
add_card(slide7, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.8), 
         "测试环境", [
             "• GPU: AMD Radeon RX 7900 XTX (24GB VRAM)",
             "• OS: Ubuntu 22.04 LTS",
             "• ROCm: 6.1",
             "• PyTorch: 2.4.0+rocm6.1",
         ], TECH_CYAN)

# 性能指标
add_card(slide7, Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.8),
         "性能指标", [
             "• 单据分类准确率: 98%+",
             "• 字段提取准确率: 95%+",
             "• 三单匹配准确率: 99%+",
             "• 单张处理时间: < 30 秒",
             "• GPU 显存占用: < 12 GB",
         ], SUCCESS_GREEN)

# 优化技术
add_card(slide7, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.2),
         "优化技术", [
             "• FP16 混合精度推理 — 减少显存占用，提升计算效率",
             "• vLLM 批处理优化 — 支持并发请求，提高吞吐量",
             "• BGE 向量化 — 本地 Embedding 计算，零网络延迟",
             "• ROCm HIP 后端 — 原生 AMD GPU 支持，无需 CUDA",
         ], AMD_RED)

# ==================== 幻灯片8：三单校验 ====================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide8)
add_title(slide8, "三单交叉校验逻辑")
add_accent_line(slide8)

# 校验规则
rules = [
    ("PO号匹配", "po_number ↔ po_reference"),
    ("物料编码匹配", "item_code ↔ item_code"),
    ("数量匹配", "quantity ↔ quantity_delivered (±5%容差)"),
    ("单价匹配", "unit_price ↔ unit_price"),
    ("金额匹配", "total_amount ↔ amount (±1%容差)"),
]

add_card(slide8, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.5),
         "校验规则", [f"• {rule}: {desc}" for rule, desc in rules], TECH_CYAN)

# 校验流程
flow_steps = [
    "1. 上传采购订单 → Extract Agent 提取字段",
    "2. 上传送货单 → Extract Agent 提取字段",
    "3. 上传发票 → Extract Agent 提取字段",
    "4. 三单交叉校验 → Validate Agent 比对",
    "5. 生成校验报告 → 通过/异常 分流处理",
]

add_card(slide8, Inches(6.8), Inches(1.5), Inches(5.7), Inches(3.5),
         "校验流程", flow_steps, SUCCESS_GREEN)

# 底部说明
add_card(slide8, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5),
         "容差设置", [
             "• 数量容差: 5% — 允许合理的发货差异",
             "• 金额容差: 1% — 严格的财务控制",
             "• 超出容差自动标记为异常，推送给审批人处理",
         ], WARNING_ORANGE)

# ==================== 幻灯片9：商业价值 ====================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide9)
add_title(slide9, "商业价值")
add_accent_line(slide9)

# 价值数据
value_data = [
    ("15 min → 2 min", "单据处理时间", "87%", AMD_RED),
    ("3% → 0.5%", "错误率", "83%", WARNING_ORANGE),
    ("8人 → 2人", "人力需求", "75%", TECH_CYAN),
    ("T+1 → 实时", "库存差异发现", "100%", SUCCESS_GREEN),
]

for i, (before, label, improvement, color) in enumerate(value_data):
    x = Inches(0.8 + i * 3.1)
    
    # 卡片背景
    card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(2.8), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = MID_BLUE
    card.line.fill.background()
    
    # 顶部颜色条
    top_bar = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.8), Inches(2.8), Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = color
    top_bar.line.fill.background()
    
    # 改进数据
    tb_before = slide9.shapes.add_textbox(x + Inches(0.2), Inches(2.2), Inches(2.4), Inches(0.8))
    tf = tb_before.text_frame
    p = tf.paragraphs[0]
    p.text = before
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(20)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE
    p.runs[0].font.name = FONT_TITLE
    
    # 标签
    tb_label = slide9.shapes.add_textbox(x + Inches(0.2), Inches(3), Inches(2.4), Inches(0.5))
    tf = tb_label.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(14)
    p.runs[0].font.color.rgb = LIGHT_GRAY
    p.runs[0].font.name = FONT_BODY
    
    # 提升百分比
    tb_imp = slide9.shapes.add_textbox(x + Inches(0.2), Inches(3.5), Inches(2.4), Inches(1))
    tf = tb_imp.text_frame
    p = tf.paragraphs[0]
    p.text = f"+{improvement}"
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(48)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = color
    p.runs[0].font.name = FONT_TITLE

# ==================== 幻灯片10：创新点 ====================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide10)
add_title(slide10, "创新点")
add_accent_line(slide10)

innovations = [
    ("领域专精 Agent", ["针对供应链单据场景深度优化", "非通用 Agent，专业能力更强"], TECH_CYAN),
    ("多 Agent 协作", ["5 个专业 Agent 各司其职", "协同完成复杂业务流程"], SUCCESS_GREEN),
    ("三单交叉校验", ["自动核对 PO/送货单/发票", "替代人工经验判断"], AMD_RED),
    ("完全本地化", ["商业数据零泄露风险", "适合企业私有化部署"], WARNING_ORANGE),
]

for i, (title, lines, color) in enumerate(innovations):
    x = Inches(0.8 + i * 3.1)
    add_card(slide10, x, Inches(1.8), Inches(2.8), Inches(2.5), title, lines, color)

# ==================== 幻灯片11：未来扩展 ====================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide11)
add_title(slide11, "未来扩展")
add_accent_line(slide11)

expansions = [
    ("图像单据支持", ["集成多模态 LLM", "支持拍照识别手写单据"], TECH_CYAN),
    ("全流程覆盖", ["扩展到出库、对账、结算", "覆盖供应链全链路"], SUCCESS_GREEN),
    ("ERP 集成", ["对接 SAP/Oracle/用友", "实现真正的自动化"], AMD_RED),
    ("多仓库协同", ["支持多仓库并行处理", "统一管理和调度"], WARNING_ORANGE),
]

for i, (title, lines, color) in enumerate(expansions):
    x = Inches(0.8 + i * 3.1)
    add_card(slide11, x, Inches(1.8), Inches(2.8), Inches(2.2), title, lines, color)

# ==================== 幻灯片12：总结 ====================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide12)

# 顶部装饰
top_bar = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = AMD_RED
top_bar.line.fill.background()

# Thank You
tb = slide12.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(1.2))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.alignment = PP_ALIGN.CENTER
p.runs[0].font.size = Pt(60)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = WHITE
p.runs[0].font.name = FONT_TITLE

# 关键特性
features = [
    "✓ Multi-Agent collaboration (5 specialized agents)",
    "✓ AMD GPU acceleration with ROCm 6.x",
    "✓ Complete privacy protection (100% local)",
    "✓ Real-world business value (87% time reduction)",
    "✓ High accuracy (98%+ classification, 99%+ matching)",
]

tb_features = slide12.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(2.5))
tf = tb_features.text_frame
tf.word_wrap = True
for i, feature in enumerate(features):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = feature
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(18)
    p.runs[0].font.color.rgb = TECH_CYAN
    p.runs[0].font.name = FONT_BODY
    p.space_after = Pt(12)

# 联系信息
tb_contact = slide12.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11.3), Inches(1.2))
tf = tb_contact.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Team: cyslmsolomon"
p.alignment = PP_ALIGN.CENTER
p.runs[0].font.size = Pt(18)
p.runs[0].font.color.rgb = LIGHT_GRAY
p.runs[0].font.name = FONT_BODY

p2 = tf.add_paragraph()
p2.text = "GitHub: github.com/cyslmsolomon/Radeon-hackathon-2026-07"
p2.alignment = PP_ALIGN.CENTER
p2.runs[0].font.size = Pt(16)
p2.runs[0].font.color.rgb = MID_GRAY
p2.runs[0].font.name = FONT_MONO

# 底部装饰
bottom_bar = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.42), prs.slide_width, Inches(0.08))
bottom_bar.fill.solid()
bottom_bar.fill.fore_color.rgb = AMD_RED
bottom_bar.line.fill.background()

# 保存演示文稿
output_path = "Supply_Chain_DocAgent_Presentation_v3.pptx"
prs.save(output_path)
print(f"演示文稿已保存: {output_path}")
