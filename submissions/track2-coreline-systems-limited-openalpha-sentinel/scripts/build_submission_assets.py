#!/usr/bin/env python3
"""Build bilingual submission documents from the Markdown sources.

This script intentionally keeps presentation dependencies outside the application
runtime. Install them in the development environment with:

    python -m pip install "python-docx>=1.1,<2" "python-pptx>=1.0,<2"
"""

from __future__ import annotations

import hashlib
import math
import re
import sys
from xml.etree import ElementTree
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable, Sequence

try:
    from docx import Document
    from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn as docx_qn
    from docx.shared import Inches as DocxInches
    from docx.shared import Pt as DocxPt
    from docx.shared import RGBColor as DocxRGBColor
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - dependency guidance is the behavior.
    raise SystemExit(
        "Missing submission-asset dependencies. Run: "
        "python -m pip install 'python-docx>=1.1,<2' 'python-pptx>=1.0,<2'"
    ) from exc


ROOT = Path(__file__).resolve().parents[1]
SUBMISSION_DIR = ROOT / "docs" / "submission"
OUTPUT_DIR = SUBMISSION_DIR / "generated"
TEAM_NAME = "Coreline Systems Limited"
REPOSITORY_URL = "https://github.com/coreline-systems/AMD-AI-DevMaster-Hackathon"
VIDEO_STATUS = "Demo video: public 3:39.9 English-subtitled recording"

LANGUAGES = {
    "en": {
        "font": "Arial",
        "fallback": "Aptos",
        "report": SUBMISSION_DIR / "PROJECT_REPORT.en.md",
        "slides": SUBMISSION_DIR / "SLIDES.en.md",
        "poster": SUBMISSION_DIR / "POSTER.en.md",
        "track": "TRACK 2 · PRIVATE AI AGENT",
        "footer": "Research and education only · AMD Radeon + ROCm",
    },
    "zh-CN": {
        "font": "PingFang SC" if sys.platform == "darwin" else "Microsoft YaHei",
        "fallback": "Microsoft YaHei",
        "report": SUBMISSION_DIR / "PROJECT_REPORT.zh-CN.md",
        "slides": SUBMISSION_DIR / "SLIDES.zh-CN.md",
        "poster": SUBMISSION_DIR / "POSTER.zh-CN.md",
        "track": "赛道二 · 私有 AI AGENT",
        "footer": "仅供研究与教育 · AMD Radeon + ROCm",
    },
}

# Quiet, work-focused palette with distinct semantic accents.
INK = "173039"
INK_2 = "29454D"
MUTED = "64767B"
PAPER = "F4F7F5"
WHITE = "FFFFFF"
GREEN = "167A68"
GREEN_DARK = "0F6254"
GREEN_LIGHT = "DDEDE8"
CORAL = "D85D4D"
CORAL_LIGHT = "F7E3DF"
GOLD = "C79B36"
GOLD_LIGHT = "F5ECD4"
LINE = "CCD7D4"
CODE_BG = "E9EFED"


@dataclass
class Block:
    kind: str
    text: str = ""
    level: int = 0
    rows: list[list[str]] = field(default_factory=list)


@dataclass
class SlideSource:
    number: int
    title: str
    preface: list[Block]
    sections: list[tuple[str, list[Block]]]


def hex_rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def docx_rgb(value: str) -> DocxRGBColor:
    value = value.lstrip("#")
    return DocxRGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def clean_markdown(text: str) -> str:
    text = text.strip()
    text = re.sub(r"!\[([^]]*)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"\[([^]]+)\]\(([^)]+)\)", r"\1 (\2)", text)
    text = text.replace("**", "").replace("__", "")
    text = re.sub(r"(?<!\*)\*(?!\*)", "", text)
    text = text.replace("`", "")
    return text.strip()


def is_table_separator(line: str) -> bool:
    cells = [cell.strip() for cell in line.strip().strip("|").split("|")]
    return bool(cells) and all(re.fullmatch(r":?-{3,}:?", cell) for cell in cells)


def is_block_start(lines: Sequence[str], index: int) -> bool:
    line = lines[index].strip()
    if not line:
        return True
    if line.startswith(("#", "```", ">", "- ", "* ", "+ ")):
        return True
    if re.match(r"\d+\.\s+", line):
        return True
    if line == "---":
        return True
    if line.startswith("|") and index + 1 < len(lines) and is_table_separator(lines[index + 1]):
        return True
    return False


def parse_markdown(text: str) -> list[Block]:
    lines = text.splitlines()
    blocks: list[Block] = []
    index = 0
    while index < len(lines):
        raw = lines[index]
        line = raw.strip()
        if not line:
            index += 1
            continue
        if line.startswith("```"):
            language = line[3:].strip()
            index += 1
            code_lines: list[str] = []
            while index < len(lines) and not lines[index].strip().startswith("```"):
                code_lines.append(lines[index].rstrip())
                index += 1
            index += 1
            blocks.append(Block("code", "\n".join(code_lines), level=1 if language else 0))
            continue
        heading = re.match(r"^(#{1,6})\s+(.+)$", line)
        if heading:
            blocks.append(Block("heading", heading.group(2).strip(), level=len(heading.group(1))))
            index += 1
            continue
        if line.startswith("|") and index + 1 < len(lines) and is_table_separator(lines[index + 1]):
            rows: list[list[str]] = []
            while index < len(lines) and lines[index].strip().startswith("|"):
                if not is_table_separator(lines[index]):
                    rows.append([cell.strip() for cell in lines[index].strip().strip("|").split("|")])
                index += 1
            if rows:
                width = max(len(row) for row in rows)
                rows = [row + [""] * (width - len(row)) for row in rows]
                blocks.append(Block("table", rows=rows))
            continue
        bullet = re.match(r"^\s*[-*+]\s+(.+)$", raw)
        if bullet:
            blocks.append(Block("bullet", bullet.group(1).strip()))
            index += 1
            continue
        ordered = re.match(r"^\s*(\d+)\.\s+(.+)$", raw)
        if ordered:
            blocks.append(Block("number", ordered.group(2).strip(), level=int(ordered.group(1))))
            index += 1
            continue
        if line.startswith(">"):
            quote_lines: list[str] = []
            while index < len(lines) and lines[index].strip().startswith(">"):
                quote_lines.append(lines[index].strip()[1:].strip())
                index += 1
            blocks.append(Block("quote", " ".join(quote_lines)))
            continue
        if line == "---":
            blocks.append(Block("rule"))
            index += 1
            continue

        paragraph_lines = [line]
        index += 1
        while index < len(lines) and not is_block_start(lines, index):
            paragraph_lines.append(lines[index].strip())
            index += 1
        blocks.append(Block("paragraph", " ".join(paragraph_lines)))
    return blocks


def parse_slide_sources(path: Path) -> list[SlideSource]:
    text = path.read_text(encoding="utf-8")
    chunks = re.split(r"^\s*---\s*$", text, flags=re.MULTILINE)
    slides: list[SlideSource] = []
    title_re = re.compile(r"^##\s+(?:Slide\s+|第\s*)(\d+)(?:\s*[-页]\s*)(.+)$", re.MULTILINE)
    for chunk in chunks:
        match = title_re.search(chunk)
        if not match:
            continue
        number = int(match.group(1))
        title = clean_markdown(match.group(2)).lstrip("-–— ").strip()
        body = chunk[match.end() :]
        section_chunks = re.split(r"^###\s+(.+)$", body, flags=re.MULTILINE)
        preface = parse_markdown(section_chunks[0])
        sections: list[tuple[str, list[Block]]] = []
        for idx in range(1, len(section_chunks), 2):
            sections.append((clean_markdown(section_chunks[idx]), parse_markdown(section_chunks[idx + 1])))
        slides.append(SlideSource(number, title, preface, sections))
    slides.sort(key=lambda item: item.number)
    if [slide.number for slide in slides] != list(range(1, 13)):
        raise ValueError(f"Expected Markdown slides 1-12 in {path}, found {[s.number for s in slides]}")
    return slides


def inline_docx(paragraph, text: str, lang: str, *, size: float | None = None, color: str = INK) -> None:
    config = LANGUAGES[lang]
    token_re = re.compile(r"(\*\*.*?\*\*|`.*?`|\[[^]]+\]\([^)]+\))")
    for token in token_re.split(text):
        if not token:
            continue
        bold = token.startswith("**") and token.endswith("**")
        code = token.startswith("`") and token.endswith("`")
        link = re.fullmatch(r"\[([^]]+)\]\(([^)]+)\)", token)
        if link:
            value = f"{link.group(1)} ({link.group(2)})"
        elif bold or code:
            value = token[2:-2] if bold else token[1:-1]
        else:
            value = token
        run = paragraph.add_run(value)
        run.bold = bold
        run.font.name = "SFMono-Regular" if code and lang == "zh-CN" else ("Aptos Mono" if code else config["font"])
        if size is not None:
            run.font.size = DocxPt(size)
        run.font.color.rgb = docx_rgb(color)
        rpr = run._element.get_or_add_rPr()
        fonts = rpr.rFonts
        if fonts is None:
            fonts = OxmlElement("w:rFonts")
            rpr.insert(0, fonts)
        fonts.set(docx_qn("w:ascii"), run.font.name)
        fonts.set(docx_qn("w:hAnsi"), run.font.name)
        fonts.set(docx_qn("w:eastAsia"), config["font"])
        fonts.set(docx_qn("w:cs"), config["fallback"])


def set_docx_cell_fill(cell, color: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(docx_qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(docx_qn("w:fill"), color)


def set_docx_cell_border(cell, color: str = LINE, size: str = "4") -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    borders = tc_pr.first_child_found_in("w:tcBorders")
    if borders is None:
        borders = OxmlElement("w:tcBorders")
        tc_pr.append(borders)
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        tag = f"w:{edge}"
        element = borders.find(docx_qn(tag))
        if element is None:
            element = OxmlElement(tag)
            borders.append(element)
        element.set(docx_qn("w:val"), "single")
        element.set(docx_qn("w:sz"), size)
        element.set(docx_qn("w:color"), color)


def add_page_number(paragraph) -> None:
    paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = paragraph.add_run()
    begin = OxmlElement("w:fldChar")
    begin.set(docx_qn("w:fldCharType"), "begin")
    instruction = OxmlElement("w:instrText")
    instruction.set(docx_qn("xml:space"), "preserve")
    instruction.text = " PAGE "
    end = OxmlElement("w:fldChar")
    end.set(docx_qn("w:fldCharType"), "end")
    run._r.extend((begin, instruction, end))


def configure_docx(document: Document, lang: str) -> None:
    section = document.sections[0]
    section.page_width = DocxInches(8.27)
    section.page_height = DocxInches(11.69)
    section.top_margin = DocxInches(0.72)
    section.bottom_margin = DocxInches(0.7)
    section.left_margin = DocxInches(0.78)
    section.right_margin = DocxInches(0.78)

    normal = document.styles["Normal"]
    normal.font.name = LANGUAGES[lang]["font"]
    normal.font.size = DocxPt(9.4 if lang == "en" else 9.2)
    normal.font.color.rgb = docx_rgb(INK)
    normal.paragraph_format.space_after = DocxPt(5)
    normal.paragraph_format.line_spacing = 1.13

    heading_sizes = {1: 28, 2: 17, 3: 12.5, 4: 10.5}
    for level, size in heading_sizes.items():
        style = document.styles[f"Heading {level}"]
        style.font.name = LANGUAGES[lang]["font"]
        style.font.size = DocxPt(size)
        style.font.bold = True
        style.font.color.rgb = docx_rgb(GREEN_DARK if level < 3 else INK_2)
        style.paragraph_format.space_before = DocxPt(11 if level > 1 else 0)
        style.paragraph_format.space_after = DocxPt(5)
        style.paragraph_format.keep_with_next = True

    header = section.header
    table = header.add_table(rows=1, cols=2, width=DocxInches(6.7))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    left, right = table.rows[0].cells
    left.text = "OPENALPHA SENTINEL"
    right.text = LANGUAGES[lang]["track"]
    right.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.RIGHT
    for cell in (left, right):
        set_docx_cell_fill(cell, INK)
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.font.name = LANGUAGES[lang]["font"]
                run.font.size = DocxPt(7.5)
                run.font.bold = True
                run.font.color.rgb = docx_rgb(WHITE)

    footer = section.footer
    p = footer.paragraphs[0]
    inline_docx(p, LANGUAGES[lang]["footer"], lang, size=7.2, color=MUTED)
    p.add_run("    ")
    add_page_number(p)


def add_docx_table(document: Document, rows: list[list[str]], lang: str) -> None:
    if not rows:
        return
    table = document.add_table(rows=len(rows), cols=len(rows[0]))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = True
    compact_size = 6.2 if len(rows[0]) >= 7 else (7.3 if len(rows[0]) >= 4 else 8.3)
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            set_docx_cell_border(cell)
            set_docx_cell_fill(cell, INK if row_index == 0 else (WHITE if row_index % 2 else PAPER))
            paragraph = cell.paragraphs[0]
            paragraph.paragraph_format.space_after = DocxPt(1.5)
            paragraph.paragraph_format.space_before = DocxPt(1.5)
            inline_docx(
                paragraph,
                value,
                lang,
                size=compact_size,
                color=WHITE if row_index == 0 else INK,
            )
            for run in paragraph.runs:
                run.bold = row_index == 0
    document.add_paragraph().paragraph_format.space_after = DocxPt(1)


def build_report(lang: str, output: Path) -> None:
    source = Path(LANGUAGES[lang]["report"])
    blocks = parse_markdown(source.read_text(encoding="utf-8"))
    document = Document()
    configure_docx(document, lang)
    document.core_properties.title = "OpenAlpha Sentinel - Project Report"
    document.core_properties.subject = "AMD AI DevMaster Hackathon Track 2"
    document.core_properties.author = TEAM_NAME

    first_h1 = True
    first_h2 = True
    first_numbered_section = True
    for block in blocks:
        if block.kind == "heading":
            text = clean_markdown(block.text)
            if block.level == 1 and first_h1:
                paragraph = document.add_paragraph()
                paragraph.paragraph_format.space_before = DocxPt(54)
                paragraph.paragraph_format.space_after = DocxPt(12)
                inline_docx(paragraph, text, lang, size=29, color=INK)
                for run in paragraph.runs:
                    run.bold = True
                first_h1 = False
                continue
            if block.level == 2 and first_h2:
                paragraph = document.add_paragraph()
                inline_docx(paragraph, text, lang, size=15, color=GREEN)
                paragraph.paragraph_format.space_after = DocxPt(20)
                first_h2 = False
                continue
            page_break_before = block.level == 2 and bool(re.match(r"\d+\.", text)) and first_numbered_section
            if page_break_before:
                first_numbered_section = False
            level = min(max(block.level, 1), 4)
            paragraph = document.add_paragraph(style=f"Heading {level}")
            if page_break_before:
                paragraph.paragraph_format.page_break_before = True
            inline_docx(paragraph, text, lang)
        elif block.kind == "paragraph":
            paragraph = document.add_paragraph()
            inline_docx(paragraph, block.text, lang)
        elif block.kind == "bullet":
            paragraph = document.add_paragraph(style="List Bullet")
            inline_docx(paragraph, block.text, lang)
            paragraph.paragraph_format.left_indent = DocxInches(0.24)
            paragraph.paragraph_format.first_line_indent = DocxInches(-0.13)
        elif block.kind == "number":
            paragraph = document.add_paragraph(style="List Number")
            inline_docx(paragraph, block.text, lang)
            paragraph.paragraph_format.left_indent = DocxInches(0.27)
            paragraph.paragraph_format.first_line_indent = DocxInches(-0.15)
        elif block.kind == "quote":
            table = document.add_table(rows=1, cols=2)
            table.columns[0].width = DocxInches(0.08)
            table.columns[1].width = DocxInches(6.4)
            set_docx_cell_fill(table.cell(0, 0), CORAL)
            set_docx_cell_fill(table.cell(0, 1), CORAL_LIGHT)
            p = table.cell(0, 1).paragraphs[0]
            p.paragraph_format.space_before = DocxPt(4)
            p.paragraph_format.space_after = DocxPt(4)
            inline_docx(p, block.text, lang, size=8.4, color=INK_2)
        elif block.kind == "table":
            add_docx_table(document, block.rows, lang)
        elif block.kind == "code":
            paragraph = document.add_paragraph()
            ppr = paragraph._p.get_or_add_pPr()
            shading = OxmlElement("w:shd")
            shading.set(docx_qn("w:fill"), CODE_BG)
            ppr.append(shading)
            paragraph.paragraph_format.left_indent = DocxInches(0.16)
            paragraph.paragraph_format.right_indent = DocxInches(0.16)
            paragraph.paragraph_format.space_before = DocxPt(4)
            paragraph.paragraph_format.space_after = DocxPt(6)
            run = paragraph.add_run(block.text)
            run.font.name = "Aptos Mono" if lang == "en" else "SFMono-Regular"
            run.font.size = DocxPt(6.8)
            run.font.color.rgb = docx_rgb(INK_2)
        elif block.kind == "rule":
            paragraph = document.add_paragraph()
            ppr = paragraph._p.get_or_add_pPr()
            borders = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(docx_qn("w:val"), "single")
            bottom.set(docx_qn("w:sz"), "8")
            bottom.set(docx_qn("w:color"), LINE)
            borders.append(bottom)
            ppr.append(borders)

    document.save(output)


class TextAudit:
    def __init__(self, name: str):
        self.name = name
        self.issues: list[str] = []

    @staticmethod
    def estimate(text: str, width: float, height: float, size: float) -> tuple[float, float]:
        if not text.strip():
            return 0.0, 0.0
        capacity = max(1.0, width * 72 / (size * 0.58))
        lines = 0
        for physical_line in text.splitlines() or [text]:
            units = sum(1.0 if ord(char) > 255 else 0.56 for char in physical_line) or 1.0
            lines += max(1, math.ceil(units / capacity))
        # The caller already subtracts text-frame margins from width/height.
        needed = lines * size * 1.22 / 72 + 0.01
        return needed, height

    def record(self, text: str, width: float, height: float, size: float, label: str) -> None:
        needed, available = self.estimate(text, width, height, size)
        if needed > available * 1.12:
            self.issues.append(f"{label}: estimated {needed:.2f}in > {available:.2f}in")

    def assert_clean(self) -> None:
        if self.issues:
            raise ValueError(f"Text overflow audit failed for {self.name}:\n" + "\n".join(self.issues))


def best_font_size(text: str, width: float, height: float, maximum: float, minimum: float) -> float:
    size = maximum
    while size > minimum:
        needed, _ = TextAudit.estimate(text, width, height, size)
        if needed <= height * 0.96:
            break
        size -= 0.5
    return max(size, minimum)


def style_ppt_run(run, lang: str, size: float, color: str, *, bold: bool = False) -> None:
    run.font.name = LANGUAGES[lang]["font"]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = hex_rgb(color)


def set_shape_fill(shape, color: str, line: str | None = None) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(color)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_rgb(line)
        shape.line.width = Pt(0.8)


def add_rect(slide, x: float, y: float, w: float, h: float, color: str, line: str | None = None, radius: bool = False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_fill(shape, color, line)
    return shape


def add_text(
    slide,
    audit: TextAudit,
    lang: str,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 18,
    minimum: float = 11,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.05,
    label: str = "text",
):
    text = clean_markdown(text)
    actual_size = best_font_size(text, max(0.2, w - margin * 2), max(0.2, h - margin * 2), size, minimum)
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = 1.02
    run = paragraph.add_run()
    run.text = text
    style_ppt_run(run, lang, actual_size, color, bold=bold)
    audit.record(text, w - margin * 2, h - margin * 2, actual_size, label)
    return shape


def add_bullets(
    slide,
    audit: TextAudit,
    lang: str,
    x: float,
    y: float,
    w: float,
    h: float,
    items: Sequence[str],
    *,
    size: float = 17,
    minimum: float = 11,
    color: str = INK,
    accent: str = GREEN,
    numbered: bool = False,
    label: str = "bullets",
):
    cleaned = [clean_markdown(item) for item in items if clean_markdown(item)]
    display = "\n".join(f"{index + 1}. {item}" if numbered else f"• {item}" for index, item in enumerate(cleaned))
    actual_size = best_font_size(display, w - 0.16, h - 0.12, size, minimum)
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    for index, item in enumerate(cleaned):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(8 if len(cleaned) <= 5 else 5)
        paragraph.line_spacing = 1.02
        lead = f"{index + 1}. " if numbered else "• "
        run = paragraph.add_run()
        run.text = lead
        style_ppt_run(run, lang, actual_size, accent, bold=True)
        run = paragraph.add_run()
        run.text = item
        style_ppt_run(run, lang, actual_size, color)
    audit.record(display, w - 0.16, h - 0.12, actual_size, label)
    return shape


def add_slide_chrome(slide, audit: TextAudit, lang: str, number: int, title: str) -> None:
    add_rect(slide, 0, 0, 0.16, 7.5, GREEN)
    add_text(slide, audit, lang, 0.48, 0.23, 2.7, 0.3, "OPENALPHA SENTINEL", size=10, minimum=10, color=GREEN_DARK, bold=True, margin=0, label="brand")
    add_text(slide, audit, lang, 10.2, 0.23, 2.55, 0.3, LANGUAGES[lang]["track"], size=9, minimum=8, color=MUTED, align=PP_ALIGN.RIGHT, margin=0, label="track")
    add_text(slide, audit, lang, 0.48, 0.68, 11.9, 0.62, title, size=28, minimum=23, color=INK, bold=True, label=f"slide {number} title")
    add_rect(slide, 0.48, 1.37, 12.25, 0.015, LINE)
    add_text(slide, audit, lang, 0.48, 7.13, 6.5, 0.22, LANGUAGES[lang]["footer"], size=7.6, minimum=7, color=MUTED, margin=0, label="footer")
    add_text(slide, audit, lang, 12.1, 7.1, 0.55, 0.26, f"{number:02d}", size=9, minimum=8, color=GREEN, bold=True, align=PP_ALIGN.RIGHT, margin=0, label="page")


def meaningful_sections(source: SlideSource, lang: str) -> list[tuple[str, list[Block]]]:
    excluded = {"Visual", "Speaker note", "视觉内容", "讲解备注"}
    return [(heading, blocks) for heading, blocks in source.sections if heading not in excluded]


def block_items(blocks: Iterable[Block], kinds: set[str] | None = None) -> list[str]:
    kinds = kinds or {"bullet", "number"}
    return [block.text for block in blocks if block.kind in kinds]


def block_paragraphs(blocks: Iterable[Block]) -> list[str]:
    return [block.text for block in blocks if block.kind in {"paragraph", "quote"}]


def first_table(blocks: Iterable[Block]) -> list[list[str]]:
    return next((block.rows for block in blocks if block.kind == "table"), [])


def first_code(blocks: Iterable[Block]) -> str:
    return next((block.text for block in blocks if block.kind == "code"), "")


def add_ppt_table(
    slide,
    audit: TextAudit,
    lang: str,
    x: float,
    y: float,
    w: float,
    h: float,
    rows: list[list[str]],
    *,
    font_size: float = 12,
    label: str = "table",
):
    if not rows:
        return None
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.035)
            cell.margin_bottom = Inches(0.035)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_rgb(INK if row_index == 0 else (WHITE if row_index % 2 else PAPER))
            frame = cell.text_frame
            frame.clear()
            paragraph = frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            run.text = clean_markdown(value)
            size = font_size - (1.0 if len(rows[0]) >= 4 else 0)
            style_ppt_run(run, lang, size, WHITE if row_index == 0 else INK, bold=row_index == 0)
            audit.record(clean_markdown(value), w / len(rows[0]) - 0.14, h / len(rows) - 0.07, size, f"{label} r{row_index}c{col_index}")
    return shape


def add_section_heading(slide, audit: TextAudit, lang: str, x: float, y: float, w: float, text: str, color: str = GREEN) -> None:
    add_rect(slide, x, y + 0.04, 0.06, 0.29, color)
    add_text(slide, audit, lang, x + 0.14, y, w - 0.14, 0.38, text, size=14, minimum=12, color=INK_2, bold=True, label="section heading")


def render_slide_1(slide, audit: TextAudit, lang: str, source: SlideSource) -> None:
    add_rect(slide, 0, 0, 13.333, 7.5, PAPER)
    add_rect(slide, 8.45, 0, 4.883, 7.5, INK)
    add_rect(slide, 0.62, 0.62, 0.72, 0.1, CORAL)
    add_text(slide, audit, lang, 0.62, 0.87, 6.95, 0.75, source.title, size=34, minimum=30, color=INK, bold=True, label="cover title")
    preface_paragraphs = block_paragraphs(source.preface)
    tagline = preface_paragraphs[0] if preface_paragraphs else ""
    add_text(slide, audit, lang, 0.62, 1.74, 6.95, 0.86, tagline, size=20, minimum=16, color=GREEN_DARK, bold=True, label="cover tagline")
    items = block_items(source.preface)
    add_bullets(slide, audit, lang, 0.65, 2.87, 6.55, 1.6, items, size=15, minimum=12, color=INK_2, accent=CORAL, label="cover metadata")
    add_text(slide, audit, lang, 0.65, 6.65, 6.7, 0.28, LANGUAGES[lang]["footer"], size=9, minimum=8, color=MUTED, margin=0, label="cover footer")

    # Abstract source-to-evidence visual.
    visual_x = 8.82
    labels = ["GitHub", "RSS / Atom", "Snapshot"]
    for idx, label in enumerate(labels):
        y = 1.25 + idx * 1.15
        add_rect(slide, visual_x, y, 1.38, 0.62, WHITE, radius=True)
        add_text(slide, audit, lang, visual_x, y + 0.05, 1.38, 0.45, label, size=11, minimum=9, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, label="source node")
        add_rect(slide, visual_x + 1.55, y + 0.29, 0.72, 0.04, GOLD)
    add_rect(slide, 11.1, 1.67, 1.72, 2.42, GREEN, radius=True)
    center = "Strategy\nCards\n+\nEvidence" if lang == "en" else "策略卡\n+\n证据"
    add_text(slide, audit, lang, 11.2, 2.08, 1.52, 1.5, center, size=17, minimum=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, label="evidence node")
    add_rect(slide, 9.1, 5.08, 3.45, 0.08, CORAL)
    local = "LOCAL RAG · CITED ANSWERS" if lang == "en" else "本地 RAG · 带引用回答"
    add_text(slide, audit, lang, 9.02, 5.35, 3.63, 0.55, local, size=13, minimum=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER, label="local claim")
    add_text(slide, audit, lang, 9.05, 6.2, 3.6, 0.5, LANGUAGES[lang]["track"], size=9, minimum=8, color="B9CAC6", align=PP_ALIGN.CENTER, label="cover track")


def render_slide_2(slide, audit: TextAudit, lang: str, source: SlideSource) -> None:
    sections = meaningful_sections(source, lang)
    heading, blocks = sections[0]
    add_section_heading(slide, audit, lang, 0.58, 1.62, 6.25, heading)
    add_bullets(slide, audit, lang, 0.66, 2.07, 6.15, 4.45, block_items(blocks), size=16, minimum=13, label="problem bullets")
    add_rect(slide, 7.35, 1.72, 5.1, 4.72, WHITE, LINE)
    source_labels = ["Code", "README", "Feeds"] if lang == "en" else ["代码", "README", "订阅源"]
    for idx, label in enumerate(source_labels):
        y = 2.12 + idx * 1.02
        color = (GREEN_LIGHT, GOLD_LIGHT, CORAL_LIGHT)[idx]
        add_rect(slide, 7.72, y, 1.34, 0.56, color)
        add_text(slide, audit, lang, 7.72, y + 0.04, 1.34, 0.4, label, size=11, minimum=9, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, label="problem source")
        add_rect(slide, 9.22, y + 0.25, 0.8, 0.035, LINE)
    add_rect(slide, 10.18, 2.3, 1.85, 2.2, INK, radius=True)
    gap = "INCONSISTENT\nEVIDENCE" if lang == "en" else "证据口径\n不一致"
    add_text(slide, audit, lang, 10.34, 2.76, 1.52, 1.15, gap, size=15, minimum=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, label="problem outcome")
    note = "Links are not a research system" if lang == "en" else "链接列表不等于研究系统"
    add_text(slide, audit, lang, 7.72, 5.23, 4.28, 0.5, note, size=14, minimum=11, color=CORAL, bold=True, align=PP_ALIGN.CENTER, label="problem note")


def render_slide_3(slide, audit: TextAudit, lang: str, source: SlideSource) -> None:
    sections = meaningful_sections(source, lang)
    heading, blocks = sections[0]
    add_section_heading(slide, audit, lang, 0.58, 1.58, 12.0, heading)
    steps = block_items(blocks)
    for idx, item in enumerate(steps[:6]):
        col = idx % 3
        row = idx // 3
        x = 0.64 + col * 4.05
        y = 2.12 + row * 1.62
        add_rect(slide, x, y, 3.62, 1.16, WHITE, LINE)
        add_rect(slide, x, y, 0.54, 1.16, GREEN if row == 0 else CORAL)
        add_text(slide, audit, lang, x + 0.08, y + 0.36, 0.38, 0.34, str(idx + 1), size=15, minimum=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER, label="step number")
        add_text(slide, audit, lang, x + 0.72, y + 0.18, 2.68, 0.78, item, size=14, minimum=11, color=INK, bold=True, valign=MSO_ANCHOR.MIDDLE, label="workflow step")
    if len(sections) > 1:
        callout = " ".join(block_paragraphs(sections[1][1]))
        add_rect(slide, 0.64, 5.64, 11.72, 0.62, GOLD_LIGHT)
        add_text(slide, audit, lang, 0.82, 5.73, 11.34, 0.4, callout, size=14, minimum=11, color=INK_2, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, label="workflow callout")


def render_slide_4(slide, audit: TextAudit, lang: str, source: SlideSource) -> None:
    sections = meaningful_sections(source, lang)
    heading, blocks = sections[0]
    add_section_heading(slide, audit, lang, 0.58, 1.58, 12.0, heading)
    items = block_items(blocks)
    abbreviations = ["OV", "SC", "AI", "SR", "JB", "WR"]
    for idx, item in enumerate(items[:6]):
        col = idx % 3
        row = idx // 3
        x = 0.64 + col * 4.05
        y = 2.15 + row * 1.75
        add_rect(slide, x, y, 3.62, 1.25, WHITE, LINE)
        add_rect(slide, x + 0.18, y + 0.2, 0.7, 0.7, (GREEN, CORAL, GOLD)[col], radius=True)
        add_text(slide, audit, lang, x + 0.18, y + 0.35, 0.7, 0.32, abbreviations[idx], size=11, minimum=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER, label="feature mark")
        add_text(slide, audit, lang, x + 1.05, y + 0.18, 2.32, 0.82, item, size=13.5, minimum=10.5, color=INK, bold=True, valign=MSO_ANCHOR.MIDDLE, label="product feature")
    note = "One local workbench · no frontend build step" if lang == "en" else "一个本地研究台 · 前端无需构建"
    add_text(slide, audit, lang, 0.66, 5.78, 11.65, 0.4, note, size=13, minimum=11, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER, label="product note")


def render_slide_5(slide, audit: TextAudit, lang: str, source: SlideSource) -> None:
    code = first_code(source.preface)
    sections = meaningful_sections(source, lang)
    add_rect(slide, 0.62, 1.7, 7.15, 4.93, INK, radius=True)
    add_text(slide, audit, lang, 0.92, 1.98, 6.55, 4.3, code, size=12.5, minimum=9.5, color=WHITE, label="architecture code")
    if sections:
        heading, blocks = sections[0]
        add_section_heading(slide, audit, lang, 8.15, 1.73, 4.45, heading, CORAL)
        add_bullets(slide, audit, lang, 8.2, 2.25, 4.15, 4.2, block_items(blocks), size=14.5, minimum=11, color=INK, accent=CORAL, label="architecture principles")


def render_slide_6(slide, audit: TextAudit, lang: str, source: SlideSource) -> None:
    rows = first_table(source.preface)
    add_ppt_table(slide, audit, lang, 0.62, 1.72, 12.05, 4.65, rows, font_size=13, label="capabilities")
    note = "5 capabilities designed · competition minimum: 2" if lang == "en" else "目标覆盖 5 项能力 · 比赛最低要求 2 项"
    add_rect(slide, 0.62, 6.52, 12.05, 0.42, GREEN_LIGHT)
    add_text(slide, audit, lang, 0.8, 6.55, 11.7, 0.28, note, size=12, minimum=10, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER, label="capability note")


def render_slide_7(slide, audit: TextAudit, lang: str, source: SlideSource) -> None:
    sections = meaningful_sections(source, lang)
    heading, blocks = sections[0]
    add_section_heading(slide, audit, lang, 0.58, 1.58, 5.8, heading)
    items = block_items(blocks)
    add_bullets(slide, audit, lang, 0.65, 2.03, 5.42, 4.45, items, size=14.5, minimum=11.2, label="evidence bullets")
    add_rect(slide, 6.55, 1.7, 5.76, 4.82, WHITE, LINE)
    add_text(slide, audit, lang, 6.88, 1.95, 4.9, 0.42, "STRATEGY CARD" if lang == "en" else "策略卡", size=14, minimum=12, color=GREEN_DARK, bold=True, label="card title")
    fields = (
        [("Source revision", "commit / immutable URL"), ("License", "SPDX + confidence"), ("Rules", "entry / exit / timeframe"), ("Disclosure", "costs / slippage / period"), ("Evidence", "field -> source fragment")]
        if lang == "en"
        else [("来源 revision", "commit / 不可变 URL"), ("许可证", "SPDX + 置信度"), ("规则", "入场 / 退出 / 周期"), ("披露", "成本 / 滑点 / 区间"), ("证据", "字段 -> 来源片段")]
    )
    for idx, (key, value) in enumerate(fields):
        y = 2.58 + idx * 0.63
        add_rect(slide, 6.9, y, 1.55, 0.42, PAPER)
        add_text(slide, audit, lang, 7.0, y + 0.04, 1.35, 0.28, key, size=10, minimum=8.5, color=MUTED, bold=True, label="card key")
        add_text(slide, audit, lang, 8.63, y + 0.02, 3.15, 0.32, value, size=11, minimum=9, color=INK, label="card value")
    if len(sections) > 1:
        callout = " ".join(block_paragraphs(sections[1][1]))
        add_rect(slide, 6.55, 6.1, 5.76, 0.54, CORAL_LIGHT)
        add_text(slide, audit, lang, 6.78, 6.18, 5.3, 0.34, callout, size=11.5, minimum=9.5, color=CORAL, bold=True, align=PP_ALIGN.CENTER, label="evidence callout")


def render_two_column_sections(slide, audit: TextAudit, lang: str, sections: list[tuple[str, list[Block]]], *, todo_right: bool = False) -> None:
    for idx, (heading, blocks) in enumerate(sections[:2]):
        x = 0.62 + idx * 6.1
        fill = GREEN_LIGHT if idx == 0 else (CORAL_LIGHT if todo_right else GOLD_LIGHT)
        accent = GREEN if idx == 0 else (CORAL if todo_right else GOLD)
        add_rect(slide, x, 1.72, 5.72, 4.95, WHITE, LINE)
        add_rect(slide, x, 1.72, 5.72, 0.58, fill)
        add_text(slide, audit, lang, x + 0.22, 1.82, 5.28, 0.34, heading, size=15, minimum=12, color=accent, bold=True, label="column heading")
        add_bullets(slide, audit, lang, x + 0.25, 2.52, 5.2, 3.77, block_items(blocks), size=15, minimum=11, color=INK, accent=accent, label="column bullets")


def render_slide_8(slide, audit: TextAudit, lang: str, source: SlideSource) -> None:
    sections = meaningful_sections(source, lang)
    content_sections = sections[:2]
    render_two_column_sections(slide, audit, lang, content_sections, todo_right=True)
    if len(sections) > 2:
        status = " ".join(block_paragraphs(sections[2][1]))
        add_text(slide, audit, lang, 0.72, 6.55, 11.75, 0.3, status, size=10.5, minimum=9, color=CORAL, bold=True, align=PP_ALIGN.CENTER, label="rocm status")


def render_slide_9(slide, audit: TextAudit, lang: str, source: SlideSource) -> None:
    sections = meaningful_sections(source, lang)
    heading, blocks = sections[0]
    add_section_heading(slide, audit, lang, 0.58, 1.55, 12.0, heading)
    rows = first_table(blocks)
    add_ppt_table(slide, audit, lang, 0.62, 2.02, 12.02, 2.55, rows, font_size=11.5, label="optimization table")
    if len(sections) > 1:
        methods = " ".join(block_paragraphs(sections[1][1]))
        add_rect(slide, 0.62, 4.88, 12.02, 0.72, GREEN_LIGHT)
        add_text(slide, audit, lang, 0.82, 4.99, 11.62, 0.47, methods, size=12.5, minimum=10, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER, label="optimization methods")
    if len(sections) > 2:
        result = " ".join(block_paragraphs(sections[2][1]))
        add_rect(slide, 0.62, 5.86, 12.02, 0.64, CORAL_LIGHT)
        add_text(slide, audit, lang, 0.82, 5.96, 11.62, 0.4, result, size=14, minimum=11, color=CORAL, bold=True, align=PP_ALIGN.CENTER, label="optimization result")


def render_slide_10(slide, audit: TextAudit, lang: str, source: SlideSource) -> None:
    sections = meaningful_sections(source, lang)
    render_two_column_sections(slide, audit, lang, sections[:2], todo_right=True)


def render_slide_11(slide, audit: TextAudit, lang: str, source: SlideSource) -> None:
    sections = meaningful_sections(source, lang)
    render_two_column_sections(slide, audit, lang, sections[:2], todo_right=True)


def render_slide_12(slide, audit: TextAudit, lang: str, source: SlideSource) -> None:
    sections = meaningful_sections(source, lang)
    if sections:
        heading, blocks = sections[0]
        add_section_heading(slide, audit, lang, 0.58, 1.55, 6.0, heading)
        add_ppt_table(slide, audit, lang, 0.62, 2.02, 6.0, 2.85, first_table(blocks), font_size=11.5, label="team table")
    if len(sections) > 1:
        heading, blocks = sections[1]
        add_section_heading(slide, audit, lang, 7.05, 1.55, 5.4, heading, CORAL)
        add_bullets(slide, audit, lang, 7.12, 2.03, 5.1, 2.95, block_items(blocks), size=13, minimum=10, accent=CORAL, label="links")
    closing = ""
    if len(sections) > 2:
        closing = " ".join(block_paragraphs(sections[2][1]))
    add_rect(slide, 0.62, 5.42, 12.02, 1.06, INK)
    add_text(slide, audit, lang, 0.9, 5.69, 11.46, 0.54, closing, size=20, minimum=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, label="closing")


SLIDE_RENDERERS = {
    1: render_slide_1,
    2: render_slide_2,
    3: render_slide_3,
    4: render_slide_4,
    5: render_slide_5,
    6: render_slide_6,
    7: render_slide_7,
    8: render_slide_8,
    9: render_slide_9,
    10: render_slide_10,
    11: render_slide_11,
    12: render_slide_12,
}


def build_slides(lang: str, output: Path) -> TextAudit:
    sources = parse_slide_sources(Path(LANGUAGES[lang]["slides"]))
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = "OpenAlpha Sentinel - Submission Slides"
    presentation.core_properties.subject = "AMD AI DevMaster Hackathon Track 2"
    presentation.core_properties.author = TEAM_NAME
    audit = TextAudit(output.name)
    blank_layout = presentation.slide_layouts[6]
    for source in sources:
        slide = presentation.slides.add_slide(blank_layout)
        add_rect(slide, 0, 0, 13.333, 7.5, PAPER)
        if source.number != 1:
            add_slide_chrome(slide, audit, lang, source.number, source.title)
        SLIDE_RENDERERS[source.number](slide, audit, lang, source)
    if len(presentation.slides) != 12:
        raise ValueError(f"Expected 12 slides, got {len(presentation.slides)}")
    audit.assert_clean()
    presentation.save(output)
    return audit


def poster_sections(path: Path) -> tuple[list[Block], dict[str, list[Block]]]:
    blocks = parse_markdown(path.read_text(encoding="utf-8"))
    preface: list[Block] = []
    sections: dict[str, list[Block]] = {}
    current = preface
    for block in blocks:
        if block.kind == "heading" and block.level == 3:
            current = []
            sections[clean_markdown(block.text)] = current
        else:
            current.append(block)
    return preface, sections


def find_poster_section(sections: dict[str, list[Block]], candidates: Sequence[str]) -> tuple[str, list[Block]]:
    for heading, blocks in sections.items():
        if any(candidate.lower() in heading.lower() for candidate in candidates):
            return heading, blocks
    return "", []


def poster_panel(slide, audit: TextAudit, lang: str, x: float, y: float, w: float, h: float, title: str, items: Sequence[str], accent: str) -> None:
    add_rect(slide, x, y, w, h, WHITE, LINE)
    add_rect(slide, x, y, w, 0.68, accent)
    add_text(slide, audit, lang, x + 0.25, y + 0.14, w - 0.5, 0.38, title, size=22, minimum=17, color=WHITE, bold=True, label="poster panel title")
    add_bullets(slide, audit, lang, x + 0.28, y + 0.93, w - 0.56, h - 1.18, items, size=18, minimum=14, color=INK, accent=accent, label="poster panel bullets")


def build_poster(lang: str, output: Path) -> TextAudit:
    source_path = Path(LANGUAGES[lang]["poster"])
    blocks = parse_markdown(source_path.read_text(encoding="utf-8"))
    _, sections = poster_sections(source_path)
    h1s = [clean_markdown(block.text) for block in blocks if block.kind == "heading" and block.level == 1]
    h2s = [clean_markdown(block.text) for block in blocks if block.kind == "heading" and block.level == 2]
    title = "OpenAlpha Sentinel"
    tagline = next((value for value in h2s if "Poster" not in value and "海报" not in value and "Export" not in value and "导出" not in value), h2s[0] if h2s else "")
    lead = next((clean_markdown(block.text) for block in blocks if block.kind == "paragraph" and "private AI agent" in block.text.lower()), "")
    if lang == "zh-CN":
        lead = next((clean_markdown(block.text) for block in blocks if block.kind == "paragraph" and "私有 AI Agent" in block.text), lead)

    presentation = Presentation()
    presentation.slide_width = Inches(33.11)
    presentation.slide_height = Inches(23.39)
    presentation.core_properties.title = "OpenAlpha Sentinel - Poster"
    presentation.core_properties.author = TEAM_NAME
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    audit = TextAudit(output.name)
    add_rect(slide, 0, 0, 33.11, 23.39, PAPER)

    # Header band.
    add_rect(slide, 0, 0, 33.11, 3.72, INK)
    add_rect(slide, 0, 0, 0.32, 3.72, CORAL)
    add_text(slide, audit, lang, 1.05, 0.55, 18.5, 1.0, title, size=46, minimum=40, color=WHITE, bold=True, label="poster title")
    add_text(slide, audit, lang, 1.08, 1.68, 20.8, 0.68, tagline, size=25, minimum=20, color="A8D7CC", bold=True, label="poster tagline")
    add_text(slide, audit, lang, 1.08, 2.48, 21.0, 0.6, lead, size=16, minimum=13, color=WHITE, label="poster lead")
    add_rect(slide, 24.4, 0.62, 7.35, 2.34, GREEN, radius=True)
    meta = f"{LANGUAGES[lang]['track']}\nAMD RADEON + ROCm\n{TEAM_NAME}"
    add_text(slide, audit, lang, 24.72, 0.93, 6.7, 1.68, meta, size=20, minimum=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, label="poster metadata")

    # Left column: problem and value.
    problem_heading, problem_blocks = find_poster_section(sections, ["fragmented", "高度分散"])
    value_heading, value_blocks = find_poster_section(sections, ["delivers", "提供"])
    poster_panel(slide, audit, lang, 0.82, 4.22, 9.72, 6.35, problem_heading, block_items(problem_blocks), CORAL)
    poster_panel(slide, audit, lang, 0.82, 10.95, 9.72, 7.03, value_heading, block_items(value_blocks), GREEN)
    boundary = next((clean_markdown(block.text) for block in value_blocks if block.kind == "paragraph"), "")
    add_rect(slide, 0.82, 18.38, 9.72, 1.42, GOLD_LIGHT)
    add_text(slide, audit, lang, 1.12, 18.66, 9.12, 0.84, boundary, size=16, minimum=13, color=INK_2, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, label="poster boundary")

    # Center column: workflow, evidence and a product mockup.
    workflow_title = "Evidence-first workflow" if lang == "en" else "证据优先工作流"
    add_section_heading(slide, audit, lang, 11.15, 4.25, 10.25, workflow_title)
    steps = ["Discover", "Normalize", "Extract", "Index", "Ask", "Watch"] if lang == "en" else ["发现", "规范化", "抽取", "索引", "问答", "监控"]
    for idx, step in enumerate(steps):
        col = idx % 3
        row = idx // 3
        x = 11.2 + col * 3.45
        y = 5.03 + row * 1.42
        add_rect(slide, x, y, 2.85, 0.92, (GREEN, GOLD, CORAL)[col], radius=True)
        add_text(slide, audit, lang, x + 0.12, y + 0.2, 2.61, 0.46, step, size=17, minimum=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, label="poster workflow")

    card_heading, card_blocks = find_poster_section(sections, ["StrategyCard preserves", "StrategyCard 保留"])
    add_section_heading(slide, audit, lang, 11.15, 8.18, 10.25, card_heading, GOLD)
    add_bullets(slide, audit, lang, 11.25, 8.85, 10.0, 4.62, block_items(card_blocks), size=16.5, minimum=13.5, color=INK, accent=GOLD, label="poster evidence")

    # Drawn workbench proof instead of an unfilled screenshot placeholder.
    mock_title = "Cited research answer" if lang == "en" else "带引用研究回答"
    add_rect(slide, 11.18, 13.72, 10.18, 6.08, WHITE, LINE)
    add_rect(slide, 11.18, 13.72, 10.18, 0.7, INK)
    add_text(slide, audit, lang, 11.55, 13.89, 7.2, 0.35, mock_title, size=18, minimum=15, color=WHITE, bold=True, label="mock title")
    add_rect(slide, 11.55, 14.78, 7.2, 0.82, PAPER, radius=True)
    question = "Which strategies disclose costs, and where are they from?" if lang == "en" else "哪些策略披露了成本？它们出自哪里？"
    add_text(slide, audit, lang, 11.82, 14.98, 6.66, 0.4, question, size=14.5, minimum=12, color=INK, label="mock question")
    answers = (
        ["Daily Equity Mean Reversion [S1]", "ETF Pairs Research [S2]", "Unknown fields remain explicit"]
        if lang == "en"
        else ["日线股票均值回归 [S1]", "ETF 配对研究 [S2]", "无证据字段保持 unknown"]
    )
    add_bullets(slide, audit, lang, 11.65, 16.0, 6.95, 2.82, answers, size=14.5, minimum=12, accent=GREEN, label="mock answer")
    add_rect(slide, 19.05, 14.8, 1.75, 3.98, GREEN_LIGHT)
    add_text(slide, audit, lang, 19.18, 15.12, 1.5, 3.2, "S1\nSOURCE\n\nS2\nSOURCE" if lang == "en" else "S1\n来源\n\nS2\n来源", size=15, minimum=12, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, label="mock citations")

    # Right column: architecture, capabilities and measured-status table.
    arch_title = "Local architecture" if lang == "en" else "本地架构"
    add_section_heading(slide, audit, lang, 22.0, 4.25, 10.25, arch_title, CORAL)
    layers = (
        ["Browser workbench", "Local API + bounded agent", "GitHub / RSS + llama.cpp", "SQLite + FTS/vector evidence"]
        if lang == "en"
        else ["浏览器研究台", "本地 API + 有界 Agent", "GitHub / RSS + llama.cpp", "SQLite + FTS/向量证据"]
    )
    for idx, layer in enumerate(layers):
        y = 5.0 + idx * 1.08
        add_rect(slide, 22.15 + idx * 0.25, y, 9.15 - idx * 0.5, 0.78, (INK, GREEN_DARK, GREEN, GOLD)[idx], radius=True)
        add_text(slide, audit, lang, 22.4 + idx * 0.25, y + 0.16, 8.65 - idx * 0.5, 0.4, layer, size=17, minimum=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER, label="poster architecture")

    capability_heading, capability_blocks = find_poster_section(sections, ["Track 2 capabilities", "赛道二能力"])
    add_section_heading(slide, audit, lang, 22.0, 9.58, 10.25, capability_heading, GREEN)
    add_bullets(slide, audit, lang, 22.12, 10.18, 9.8, 3.1, block_items(capability_blocks), size=16.5, minimum=13, accent=GREEN, label="poster capabilities")

    amd_heading, amd_blocks = find_poster_section(sections, ["AMD Radeon / ROCm"])
    add_section_heading(slide, audit, lang, 22.0, 13.55, 10.25, amd_heading, CORAL)
    amd_items = block_items(amd_blocks)
    add_bullets(slide, audit, lang, 22.12, 14.12, 9.8, 2.58, amd_items, size=15, minimum=12, accent=GREEN, label="poster amd verified")
    result_rows = next((block.rows for block in blocks if block.kind == "table" and any("Decode" in cell or "解码" in cell for row in block.rows for cell in row)), [])
    add_ppt_table(slide, audit, lang, 22.08, 16.95, 9.95, 2.85, result_rows, font_size=13, label="poster todo metrics")

    # Footer.
    add_rect(slide, 0, 20.55, 33.11, 2.84, INK)
    why_heading, why_blocks = find_poster_section(sections, ["Why it is different", "差异化价值"])
    why_paragraphs = block_paragraphs(why_blocks)
    why = why_paragraphs[0] if why_paragraphs else ""
    add_text(slide, audit, lang, 0.95, 20.92, 23.0, 0.52, why_heading, size=19, minimum=16, color=GOLD, bold=True, label="poster why heading")
    add_text(slide, audit, lang, 0.95, 21.52, 23.0, 0.9, why, size=16, minimum=13, color=WHITE, bold=True, label="poster why")
    footer_links = f"{REPOSITORY_URL}\n{VIDEO_STATUS}\n{TEAM_NAME}"
    add_text(slide, audit, lang, 25.0, 20.95, 7.0, 1.55, footer_links, size=16, minimum=11, color="A8D7CC", bold=True, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE, label="poster links")

    audit.assert_clean()
    presentation.save(output)
    return audit


def verify_docx(path: Path) -> None:
    if not path.is_file() or path.stat().st_size < 30_000:
        raise ValueError(f"DOCX is missing or unexpectedly small: {path}")
    document = Document(path)
    if len(document.paragraphs) < 80:
        raise ValueError(f"DOCX content self-check failed: {path} has {len(document.paragraphs)} paragraphs")
    if len(document.tables) < 8:
        raise ValueError(f"DOCX table self-check failed: {path} has {len(document.tables)} tables")


def verify_pptx(path: Path, expected_slides: int) -> None:
    if not path.is_file() or path.stat().st_size < 25_000:
        raise ValueError(f"PPTX is missing or unexpectedly small: {path}")
    presentation = Presentation(path)
    if len(presentation.slides) != expected_slides:
        raise ValueError(f"Expected {expected_slides} slides in {path}, got {len(presentation.slides)}")
    if not all(any(getattr(shape, "has_text_frame", False) and shape.text.strip() for shape in slide.shapes) for slide in presentation.slides):
        raise ValueError(f"At least one slide contains no text: {path}")


def build_all() -> list[Path]:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    outputs: list[Path] = []
    for lang in LANGUAGES:
        report = OUTPUT_DIR / f"OpenAlpha-Sentinel-Project-Report.{lang}.docx"
        slides = OUTPUT_DIR / f"OpenAlpha-Sentinel-Slides.{lang}.pptx"
        poster = OUTPUT_DIR / f"OpenAlpha-Sentinel-Poster.{lang}.pptx"
        build_report(lang, report)
        build_slides(lang, slides)
        build_poster(lang, poster)
        verify_docx(report)
        verify_pptx(slides, 12)
        verify_pptx(poster, 1)
        outputs.extend((report, slides, poster))
    return outputs


def write_checksum_manifest() -> Path:
    manifest = OUTPUT_DIR / "SHA256SUMS"
    lines: list[str] = []
    for path in sorted(OUTPUT_DIR.iterdir(), key=lambda item: item.name):
        if not path.is_file() or path == manifest:
            continue
        digest = hashlib.sha256(path.read_bytes()).hexdigest()
        lines.append(f"{digest}  {path.name}")
    manifest.write_text("\n".join(lines) + "\n", encoding="ascii")
    return manifest


def normalize_public_test_evidence() -> None:
    """Remove workstation hostnames from JUnit files before publishing them."""

    for path in OUTPUT_DIR.glob("*test-results*.xml"):
        tree = ElementTree.parse(path)
        public_hostname = (
            "rocm-validation" if "rocm" in path.name.lower() else "local-validation"
        )
        for suite in tree.getroot().iter("testsuite"):
            suite.set("hostname", public_hostname)
        tree.write(path, encoding="utf-8", xml_declaration=True)


def main() -> None:
    outputs = build_all()
    normalize_public_test_evidence()
    manifest = write_checksum_manifest()
    print("Generated and verified submission assets:")
    for path in outputs:
        print(f"- {path.relative_to(ROOT)} ({path.stat().st_size:,} bytes)")
    print(f"- {manifest.relative_to(ROOT)} ({manifest.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
