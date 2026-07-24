from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色主题 - AMD风格
AMD_RED = RGBColor(0xED, 0x1C, 0x24)
DARK_BG = RGBColor(0x1F, 0x1F, 0x1F)
LIGHT_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_GRAY = RGBColor(0xC8, 0xD3, 0xE6)
TEAL = RGBColor(0x0D, 0x94, 0x88)

# ========== 幻灯片1：标题页 ==========
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

# 背景
bg = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BG
bg.line.fill.background()

# AMD红色装饰条
accent = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(2.8), prs.slide_width, Inches(0.1)
)
accent.fill.solid()
accent.fill.fore_color.rgb = AMD_RED
accent.line.fill.background()

# 标题
title = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(1.5))
tf = title.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Supply Chain DocAgent"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(48)
r.font.bold = True
r.font.color.rgb = LIGHT_TEXT

# 副标题
subtitle = slide1.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.3), Inches(1))
tf = subtitle.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "供应链单据智能处理 Agent"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(28)
r.font.color.rgb = ACCENT_GRAY

# 团队信息
team = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(11.3), Inches(1))
tf = team.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "AMD AI DevMaster Hackathon 2026 | Track 2: Private AI Agent"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(18)
r.font.color.rgb = TEAL

# ========== 幻灯片2：问题与挑战 ==========
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "供应链单据处理的挑战"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# 问题列表
problems = [
    ("效率低下", "人工逐单核对，处理速度慢，难以应对大批量单据"),
    ("易出错", "手动录入容易产生数据错误，导致库存差异"),
    ("难追溯", "纸质单据难以追踪，审计困难"),
    ("成本高", "大量人工投入，人力成本居高不下"),
]

for i, (title_text, desc) in enumerate(problems):
    y_pos = Inches(1.5 + i * 1.3)
    
    # 红色装饰点
    dot = slide2.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1), y_pos + Inches(0.1), Inches(0.3), Inches(0.3)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = AMD_RED
    dot.line.fill.background()
    
    # 问题标题
    tb = slide2.shapes.add_textbox(Inches(1.6), y_pos, Inches(10), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.runs[0].font.size = Pt(24)
    p.runs[0].font.bold = True
    
    # 问题描述
    tb2 = slide2.shapes.add_textbox(Inches(1.6), y_pos + Inches(0.4), Inches(10), Inches(0.5))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.runs[0].font.size = Pt(16)
    p2.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# ========== 幻灯片3：解决方案 ==========
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "我们的解决方案"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# 解决方案框
solutions = [
    ("多Agent协作", "Classify Agent\n识别单据类型", TEAL),
    ("智能提取", "Extract Agent\n提取关键字段", RGBColor(0x12, 0x45, 0x7A)),
    ("三单校验", "Validate Agent\n交叉核对PO/送货单/发票", RGBColor(0x2C, 0x5F, 0x2D)),
    ("异常处理", "Exception Agent\n自动分类并推送审批", AMD_RED),
    ("ERP录入", "Entry Agent\n校验通过后自动录入", RGBColor(0x6D, 0x2E, 0x46)),
]

for i, (title_text, desc, color) in enumerate(solutions):
    x_pos = Inches(0.5 + i * 2.5)
    
    # 彩色框
    box = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.5), Inches(2.2), Inches(2.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    
    # 标题
    tb = slide3.shapes.add_textbox(x_pos + Inches(0.1), Inches(1.8), Inches(2), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(20)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = LIGHT_TEXT
    
    # 描述
    tb2 = slide3.shapes.add_textbox(x_pos + Inches(0.1), Inches(2.6), Inches(2), Inches(1.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.alignment = PP_ALIGN.CENTER
    p2.runs[0].font.size = Pt(14)
    p2.runs[0].font.color.rgb = LIGHT_TEXT

# 底部说明
note = slide3.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(2))
tf = note.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "所有Agent在AMD Radeon GPU + ROCm上完全本地运行，确保数据隐私和处理性能"
p.alignment = PP_ALIGN.CENTER
p.runs[0].font.size = Pt(18)
p.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# ========== 幻灯片4：技术架构 ==========
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "技术架构"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# 架构图（用形状表示）
layers = [
    ("输入层", "拍照/扫描 | PDF/Excel | 邮件附件", RGBColor(0xE9, 0xDF, 0xC7)),
    ("Agent编排层", "Classify | Extract | Validate | Entry | Exception", TEAL),
    ("RAG引擎", "单据模板知识库 + 历史异常案例库", RGBColor(0x12, 0x45, 0x7A)),
    ("LLM推理", "vLLM + ROCm | AMD Radeon GPU加速", AMD_RED),
    ("输出层", "核对通过 → 自动入库 | 异常 → 推送审批", RGBColor(0x2C, 0x5F, 0x2D)),
]

for i, (title_text, desc, color) in enumerate(layers):
    y_pos = Inches(1.2 + i * 1.1)
    
    # 层框
    box = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), y_pos, Inches(11), Inches(0.9)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    
    # 标题
    tb = slide4.shapes.add_textbox(Inches(1.5), y_pos + Inches(0.1), Inches(3), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.runs[0].font.size = Pt(20)
    p.runs[0].font.bold = True
    
    # 描述
    tb2 = slide4.shapes.add_textbox(Inches(5), y_pos + Inches(0.1), Inches(6.5), Inches(0.7))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.runs[0].font.size = Pt(16)

# ========== 幻灯片5：AMD GPU优势 ==========
slide5 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "AMD Radeon GPU + ROCm 优势"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# 优势列表
advantages = [
    ("完全开源", "ROCm软件栈完全开源，无供应商锁定", "🔓"),
    ("本地部署", "数据完全本地处理，保护企业隐私", "🏠"),
    ("高性能", "AMD GPU强大的并行计算能力", "⚡"),
    ("低成本", "相比NVIDIA方案更具成本效益", "💰"),
    ("生态兼容", "支持PyTorch、vLLM等主流框架", "🔗"),
]

for i, (title_text, desc, icon) in enumerate(advantages):
    y_pos = Inches(1.3 + i * 1.1)
    
    # 图标
    tb_icon = slide5.shapes.add_textbox(Inches(1), y_pos, Inches(0.8), Inches(0.8))
    tf = tb_icon.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.runs[0].font.size = Pt(32)
    
    # 标题
    tb = slide5.shapes.add_textbox(Inches(2), y_pos, Inches(10), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.runs[0].font.size = Pt(24)
    p.runs[0].font.bold = True
    
    # 描述
    tb2 = slide5.shapes.add_textbox(Inches(2), y_pos + Inches(0.4), Inches(10), Inches(0.5))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.runs[0].font.size = Pt(16)
    p2.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# ========== 幻灯片6：核心能力 ==========
slide6 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "核心能力（对应赛事评分）"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# 能力表格
capabilities = [
    ("能力", "实现方式", "评分项"),
    ("本地知识检索（RAG）", "单据模板库 + 历史案例检索", "Agent功能完整性"),
    ("工具调用", "ERP API、邮件发送、文件操作", "工具调用与工作流编排"),
    ("多步骤任务规划", "识别→提取→校验→录入→异常处理", "多步骤任务规划"),
    ("本地多轮记忆", "对话历史 + 处理记录持久化", "本地多轮记忆"),
    ("权限控制", "按角色分权（录入员/审批人/管理员）", "隐私保护机制"),
]

for row_idx, row in enumerate(capabilities):
    y_pos = Inches(1.3 + row_idx * 0.9)
    is_header = row_idx == 0
    
    for col_idx, cell_text in enumerate(row):
        x_pos = Inches(0.5 + col_idx * 4.2)
        
        # 单元格
        box = slide6.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x_pos, y_pos, Inches(4), Inches(0.8)
        )
        
        if is_header:
            box.fill.solid()
            box.fill.fore_color.rgb = DARK_BG
        else:
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5) if row_idx % 2 == 0 else LIGHT_TEXT
        box.line.fill.background()
        
        # 文本
        tb = slide6.shapes.add_textbox(x_pos + Inches(0.2), y_pos + Inches(0.15), Inches(3.6), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cell_text
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(14) if not is_header else Pt(16)
        p.runs[0].font.bold = is_header
        p.runs[0].font.color.rgb = LIGHT_TEXT if is_header else DARK_BG

# ========== 幻灯片7：演示流程 ==========
slide7 = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
title = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "演示流程"
p.runs[0].font.size = Pt(36)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = DARK_BG

# 流程步骤
steps = [
    ("1", "系统启动", "启动Agent系统\n加载RAG知识库", TEAL),
    ("2", "上传采购订单", "自动识别PO类型\n提取物料/数量/价格", RGBColor(0x12, 0x45, 0x7A)),
    ("3", "上传送货单", "三单交叉校验\n检测数量/价格差异", RGBColor(0x2C, 0x5F, 0x2D)),
    ("4", "上传发票", "完成三单核对\n校验通过/异常处理", AMD_RED),
]

for i, (num, title_text, desc, color) in enumerate(steps):
    x_pos = Inches(0.5 + i * 3.2)
    
    # 圆圈数字
    circle = slide7.shapes.add_shape(
        MSO_SHAPE.OVAL, x_pos + Inches(0.9), Inches(1.5), Inches(1), Inches(1)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    
    # 数字
    tb_num = slide7.shapes.add_textbox(x_pos + Inches(0.9), Inches(1.6), Inches(1), Inches(0.8))
    tf = tb_num.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(36)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = LIGHT_TEXT
    
    # 标题
    tb = slide7.shapes.add_textbox(x_pos, Inches(2.8), Inches(2.8), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(20)
    p.runs[0].font.bold = True
    
    # 描述
    tb2 = slide7.shapes.add_textbox(x_pos, Inches(3.5), Inches(2.8), Inches(1.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.alignment = PP_ALIGN.CENTER
    p2.runs[0].font.size = Pt(14)
    p2.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# ========== 幻灯片8：总结 ==========
slide8 = prs.slides.add_slide(prs.slide_layouts[6])

# 背景
bg = slide8.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BG
bg.line.fill.background()

# 标题
title = slide8.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(1.5))
tf = title.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Thank You"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(56)
r.font.bold = True
r.font.color.rgb = LIGHT_TEXT

# 关键特性
features = [
    "✓ Multi-Agent collaboration",
    "✓ AMD GPU acceleration with ROCm",
    "✓ Complete privacy protection",
    "✓ Real-world business value",
]

for i, feature in enumerate(features):
    tb = slide8.shapes.add_textbox(Inches(3), Inches(3 + i * 0.6), Inches(7), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = feature
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(24)
    p.runs[0].font.color.rgb = TEAL

# 联系信息
contact = slide8.shapes.add_textbox(Inches(1), Inches(6), Inches(11.3), Inches(1))
tf = contact.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Team: cyslmsolomon | GitHub: github.com/cyslmsolomon/Radeon-hackathon-2026-07"
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(16)
r.font.color.rgb = ACCENT_GRAY

# 保存演示文稿
output_path = "Supply_Chain_DocAgent_Presentation.pptx"
prs.save(output_path)
print(f"演示文稿已保存: {output_path}")
